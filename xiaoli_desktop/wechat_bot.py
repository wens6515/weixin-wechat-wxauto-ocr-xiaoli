# -*- coding: utf-8 -*-
import time
import json
import os
import ctypes
import logging
import random
import re
import threading
import traceback
import requests
import base64
import pyautogui
# 无人值守自动化（用户定案）：小漓共享物理鼠标且处理消息时要前置微信点击，
# 用户鼠标恰停在/甩到屏幕四角时 fail-safe 会拒绝全部输入动作——回复静默
# 丢失且失败会话反复重处理白烧 API。「人工夺回控制权」由暂停/恢复（托盘/
# CLI）承担，不依赖此保护，故关闭。
pyautogui.FAILSAFE = False
import tempfile
from wx_backend import create_backend, BackendUnavailableError
from wx_backend.models import MessageType
from wx_backend.visual_backend import (
    ensure_window_visible,
    find_window_by_title,
    window_rect,
    default_right_half_rect,
    position_window_visible,
)
from xiaoli_app.config_store import AI_DEFAULTS
from xiaoli_app.usage_store import UsageStore
from xiaoli_app.web_search import (web_search, web_fetch, WebSearchError,
                                   format_search_results)


def models_endpoint(chat_url):
    """由聊天端点推导模型列表端点：把末尾的 /chat/completions 换成 /models。
    用 rsplit 只替换最后一处（str.replace 会替换所有出现处——自定义端点
    URL 里同一子串出现多次时拼接错误）；不含该子串时原样返回（保持旧行为，
    请求是否有效由调用方/API 决定）。"""
    if not chat_url or "/chat/completions" not in chat_url:
        return chat_url
    return chat_url.rsplit("/chat/completions", 1)[0] + "/models"


def is_group_chat(chat_name, title=None):
    """群聊判定（集中判定点）。

    title 优先：右侧会话标题形如 '强盗"集团(5)'（括号内人数，真机标定）——
    含括号人数 = 群聊。这是权威信号，普通群名（如'哆菈A夢'）不含'群/集团'
    字，名称启发式会漏判。
    title 缺省（无视觉后端/旧路径/单测）回退名称启发式：含「群」或「集团」。
    """
    if title is not None:
        t = (title or "").strip()
        return bool(re.match(r"^.+?\(\d+\)\s*$", t))
    name = chat_name or ""
    return "群" in name or "集团" in name


def strip_model_prefix(model):
    """剥离模型名的厂商前缀：'deepseek:deepseek-v4-flash' → 'deepseek-v4-flash'。

    配置/卡片里模型 id 沿用「厂商:模型」前缀格式（UI 分组展示用），
    但 API 只认纯模型名——透传带前缀的 model 会得到 400。
    """
    if not model or ":" not in model:
        return model
    return model.split(":", 1)[-1]


def estimate_tokens(text):
    """粗略估算 token 数：中英混合加权，整体偏保守上界。

    原实现 1 字符 = 1 token：对英文高估约 4 倍（实际 ~4 字符/token），
    长英文文档被过度裁剪（100k 预算实际只用 ~25k，上下文利用率低）。
    加权：CJK 0.8 token/字（实际 ~0.6，保守）、ASCII 0.3（实际 ~0.25）、
    其余 0.6。混合场景仍略偏保守（宁可多估不超限——请求超限会得到
    API 400 "maximum context length"，低估反而更危险）。
    """
    if not text:
        return 0
    cjk = ascii_n = other = 0
    for ch in text:
        o = ord(ch)
        if 0x4E00 <= o <= 0x9FFF:
            cjk += 1
        elif o < 128:
            ascii_n += 1
        else:
            other += 1
    return int(cjk * 0.8 + ascii_n * 0.3 + other * 0.6) + 1


def _content_to_text(content):
    """把消息 content 转成纯文本估算串（仅供 estimate_tokens 估算用，
    不修改原消息）。多模态块列表（vision user 消息）取 text 块的文本
    拼接，image_url 等非文本块按固定占位计——base64 字符数不代表
    token 数（图片 token 由视觉模型内部处理），按字符算会误判超预算。
    """
    if isinstance(content, list):
        parts = []
        for b in content:
            if isinstance(b, dict) and b.get("type") == "text":
                parts.append(str(b.get("text") or ""))
            else:
                parts.append("[image]")
        return "".join(parts)
    return str(content or "")


def _trim_blocks(blocks, keep):
    """多模态块列表按字符预算裁剪：text 块截断文本，image_url 等非文本
    块必须保留（截断 base64 会损坏图片；视觉调用图片是识别对象）。
    预算优先分配给非文本块（每块按固定成本计），text 块共享剩余预算
    按序截断。返回裁剪后的块列表。
    """
    fixed = sum(500 for b in blocks
                if not (isinstance(b, dict) and b.get("type") == "text"))
    text_budget = max(0, keep - fixed)
    out = []
    used_text = 0
    for b in blocks:
        if isinstance(b, dict) and b.get("type") == "text":
            text = str(b.get("text") or "")
            room = text_budget - used_text
            if room <= 0:
                out.append({"type": "text", "text": "[内容过长已截断]…"})
                continue
            if len(text) > room:
                out.append({"type": "text",
                            "text": text[:room] + "[内容过长已截断]…"})
                used_text = text_budget
            else:
                out.append(b)
                used_text += len(text)
        else:
            # 非文本块（image_url 等）：token 大头，原样保留
            out.append(b)
    return out


def fit_messages_in_budget(messages, budget=100000, reserve=2000):
    """把 messages 裁剪到 token 预算内（从最旧的非 system 消息开始丢弃）。

    根因：文件识别把超大文件全文拼进 prompt（实测请求 272 万 token，
    模型上限 104 万 → API 400 "maximum context length"）。
    规则（保序——消息布局是缓存优化的一部分，中途的 system（重要记忆/
    相关记忆/当前时间）绝不能被搬到最前）：
    - system 消息永不丢弃；单条超预算 20% 时截断其内容（人设可能很长，
      其余 system 都很短不受影响）
    - 非 system 消息保持相对顺序，从最旧（列表前端）开始丢弃直到 ≤ 预算
    - 末条 user 消息（当前消息）不丢弃，仍超预算时截断内容到预算 60%
    返回裁剪后的 messages。
    """
    budget = max(1000, budget)
    cap = max(500, budget - reserve)
    trimmed = []
    for m in messages:
        if m.get("role") == "system":
            content = str(m.get("content") or "")
            limit = max(500, budget // 5)
            if estimate_tokens(content) > limit:
                # 保留开头（system prompt 语义在前）
                trimmed.append({"role": "system", "content": content[:limit]})
            else:
                trimmed.append(m)
        else:
            trimmed.append(m)
    total = sum(estimate_tokens(_content_to_text(m.get("content")))
                for m in trimmed)
    last_idx = len(trimmed) - 1
    for idx in range(last_idx):
        if total <= cap:
            break
        m = trimmed[idx]
        if m.get("role") == "system":
            continue
        total -= estimate_tokens(_content_to_text(m.get("content")))
        trimmed[idx] = None
    out = [m for m in trimmed if m is not None]
    # 末条 user（当前消息，可能含文件全文/多模态块）仍超预算 → 截断到 60%
    if out and out[-1].get("role") == "user":
        keep = max(500, int(budget * 0.6))
        raw = out[-1].get("content")
        if estimate_tokens(_content_to_text(raw)) > keep:
            if isinstance(raw, list):
                content = _trim_blocks(raw, keep)
            else:
                content = str(raw or "")[:keep] + "[内容过长已截断]…"
            out[-1] = dict(out[-1], content=content)
    return out

LOG_FILE = os.path.join(os.path.dirname(os.path.abspath(__file__)), "bot.log")
# 前端视图日志（INFO 及以上）：GUI 日志页/首页运行日志区只读这个文件，
# 事件流干净；排障看 bot.log 全量（含 DEBUG）。
RUN_LOG_FILE = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                            "bot_run.log")

# 日志双轨轮转：bot.log 全量 DEBUG、bot_run.log 前端 INFO，各 2MB 轮转
# 保留 2 份历史。历史缺陷：模块级 import 时即清空 bot.log——每次启动丢
# 日志（排障无法追溯），且 GUI/CLI 并发 import 互清。轮转后 LogPage 的
# size<offset 检测会自动重置增量读取位置，无需改动。
from logging.handlers import RotatingFileHandler

_file_handler = RotatingFileHandler(LOG_FILE, maxBytes=2 * 1024 * 1024,
                                    backupCount=2, encoding="utf-8")
_run_handler = RotatingFileHandler(RUN_LOG_FILE, maxBytes=2 * 1024 * 1024,
                                   backupCount=2, encoding="utf-8")
_run_handler.setLevel(logging.INFO)
logging.basicConfig(
    level=logging.DEBUG,
    format='%(asctime)s [%(levelname)s] %(message)s',
    handlers=[_file_handler, _run_handler, logging.StreamHandler()],
    force=True
)
# 终端只显示 INFO 及以上（handlers[2] = StreamHandler），bot.log 保留 DEBUG
logging.getLogger().handlers[2].setLevel(logging.INFO)
logger = logging.getLogger("xiaoli")


# chat_model 为空时的共用兜底（chat / vision 两链路）：DeepSeek vision-exp。
# 必须是 strip 后的纯名——API 只认纯模型名，带厂商前缀的兜底请求必 400
# （历史缺陷：旧常量带 "deepseek:" 前缀，活跃卡缺失被模板补建（chat_model
# 空）后 vision 兜底请求全部 400）。call_chat_ai 同样兜底，防空模型原样发出。
VISION_MODEL_DEFAULT = "deepseek-v4-flash-vision-exp"

# ---------- LLM API 调用韧性（chat / vision 共用，见 _post_chat_completions） ----------

RETRY_AFTER_GIVEUP = 15.0     # 服务端 Retry-After 超过该秒数视为长时限流，放弃重试
BACKOFF_BASE = 1.0            # 指数退避起始秒数：1s → 2s → 4s …
BACKOFF_CAP = 8.0             # 单次退避封顶
API_WALL_BUDGET_DEFAULT = 45  # 墙钟预算默认值：重试总时长封顶，杜绝「超时×重试」叠成分钟级等待

# vision 工具循环的补全调用上限：搜索 + 抓取（含换源重试一次）+ 作答（防循环烧钱）
VISION_TOOL_ROUNDS = 4


class ApiCallError(Exception):
    """LLM API 调用最终失败（4xx 不可重试 / Retry-After 过长 / 重试耗尽 / 预算用尽）。

    status 保留 HTTP 状态码（网络异常时为 None）。调用方负责降级或转
    友好文案，绝不把异常文本原样发给微信好友。
    """

    def __init__(self, message, status=None):
        super().__init__(message)
        self.status = status


# API 最终失败时的角色内兜底回复（不写入对话历史；池子随机避免机器人复读同句）
FRIENDLY_API_ERROR_REPLIES = (
    "呜…API 那头刚刚没了回应(´･ω･`)，不是不理你，稍后再喊我一次好不好",
    "(>_<) 后端突然开小差了，这句没接住，你再来一次我肯定在",
    "刚刚信号断了一拍(´･ω･`)，等下再发一遍吧，这次一定接住",
)


# 纯图/表情消息的 vision 单调用指令（与 xiaoli_bot.VISION_ROUTE_PROMPT 同一
# 分流语义的纯图变体：看图后角色化回复，任务材料则投递）。历史缺陷：此路径
# 曾用旧两段式时代的 vision_prompt（「专业图像描述AI，客观复述图片」），把
# 复述文本当角色回复原样发出——已废弃统一到角色化链路。
VISION_IMAGE_PROMPT = (
    "请严格遵守以上人设（包括不用 emoji、改用颜文字）。用户发来了一张或多张图片"
    "（没有附带文字；多张时按发送顺序给出，可能相互关联，请结合起来看）。"
    "请看图后直接以你的身份回复用户（可以评价、接梗、回答图里的问题）；"
    "如果图片明显是任务材料（如文档截图、带指令的截图、需要处理的内容），"
    "调用 dispatch_task 工具投递任务，任务描述写在工具参数里。"
)


# memory 键归一化（用户定案方案 A）：OCR 对引号半/全角极不稳（'强盗”集团'
# / '强盗"集团' / '" 强盗 " 集团' 是同一会话），原样做键会让同一会话的
# 记忆分裂到多个条目。规则：所有引号变体（单双/半全角/弯直）与空格类
# 字符一律剥掉再存取。只影响 memory 键——显示名（回复目标/日志/群聊
# 判定）不受影响。剥引号同时覆盖「OCR 整个丢掉前引号」的漏字场景。
_QUOTE_CHARS = (
    "\u201c\u201d\u2018\u2019\u201e\u201f"  # “ ” ‘ ’ „ ‟
    "\u00ab\u00bb\u2039\u203a"              # « » ‹ ›
    "\u300c\u300d\u300e\u300f"              # 「 」 『 』
    "\uff02\u02bc\u0060\u00b4\"'"           # ＂ ʼ ` ´ " '
)


def _memory_key(chat_id):
    """memory 键归一化：剥掉所有引号变体与空格类字符。"""
    s = str(chat_id or "").translate(str.maketrans("", "", _QUOTE_CHARS))
    return re.sub(r"\s+", "", s)


def migrate_memory_data(data):
    """memory.json 结构迁移（v2）：{chat: {"recent": [...], "important":
    [...], "index": [...], "indexed": N}}。

    旧结构 {chat: [msgs]} 原位升级——旧消息整体作为 recent（超出 recent
    上限的部分由启动加载流程一次性归档进深层文件）；坏数据项返回空。
    """
    if not isinstance(data, dict):
        return {}
    out = {}
    for chat, v in data.items():
        if isinstance(v, dict) and isinstance(v.get("recent"), list):
            out[chat] = {
                "recent": v["recent"],
                "important": [x for x in (v.get("important") or [])
                              if isinstance(x, dict)],
                "index": [x for x in (v.get("index") or [])
                          if isinstance(x, dict)],
                "indexed": int(v.get("indexed") or 0),
            }
        elif isinstance(v, list):
            out[chat] = {"recent": v, "important": [], "index": [], "indexed": 0}
    return out


# ---------- 深层文件级操作（bot 未运行时 UI 直读/直删走这里，与 bot 方法共用） ----------

def deep_count_lines(path):
    """深层 jsonl 有效行数（坏行不计）。文件缺失返回 0。"""
    n = 0
    try:
        with open(path, "r", encoding="utf-8") as f:
            for line in f:
                line = line.strip()
                if line.startswith("{"):
                    n += 1
    except OSError:
        return 0
    return n


def deep_read_page(path, offset=0, limit=200, query=None):
    """深层 jsonl 分页/过滤读取。

    query 非空：全量按内容子串（不区分大小写）过滤，返回 (命中列表截
    200 条, 命中总数)；否则返回 (第 offset 页的 ≤limit 条, None)。
    """
    msgs = []
    try:
        with open(path, "r", encoding="utf-8") as f:
            for line in f:
                line = line.strip()
                if not line:
                    continue
                try:
                    item = json.loads(line)
                except ValueError:
                    continue
                if isinstance(item, dict):
                    msgs.append(item)
    except OSError:
        return [], None
    if query:
        q = str(query).lower()
        hits = [m for m in msgs if q in str(m.get("content") or "").lower()]
        return hits[:200], len(hits)
    return msgs[offset:offset + limit], None


def deep_delete_line(path, line_no):
    """删除深层 jsonl 第 line_no（1 基，按文件序=时间正序）行。

    tmp + replace 原子重写。返回 (被删消息, 新有效行数)；行号越界/
    文件缺失/重写失败返回 (None, 原行数)。"""
    try:
        with open(path, "r", encoding="utf-8") as f:
            lines = f.readlines()
    except OSError:
        return None, 0
    kept, removed, n = [], None, 0
    for line in lines:
        stripped = line.strip()
        if not stripped or not stripped.startswith("{"):
            kept.append(line)  # 坏行原样保留
            continue
        n += 1
        if n == line_no:
            try:
                removed = json.loads(stripped)
            except ValueError:
                pass
            continue
        kept.append(line)
    if removed is None:
        return None, n
    try:
        tmp = path + ".tmp"
        with open(tmp, "w", encoding="utf-8") as f:
            f.writelines(kept)
        os.replace(tmp, path)
    except OSError as e:
        logger.error(f"[记忆] 深层文件重写失败: {e}")
        return None, n
    return removed, n - 1


def load_config(path="config.json"):
    default_cfg = {
        "bot_nickname": "小漓",
        "ai_api_url": "https://api.deepseek.com/v1/chat/completions",
        "ai_api_key": "",
        "chat_model": "deepseek:deepseek-v4-flash",
        "chat_temperature": 0.7,
        "chat_top_p": 0.9,
        "system_prompt": AI_DEFAULTS["system_prompt"],
        "max_history": 1000,
        "cooldown": 3,
        "api_retry": 2,
        "api_timeout": 60,
        "api_wall_budget": API_WALL_BUDGET_DEFAULT,
        "wechat_window_rect": None,
        "start_paused": True,
        "memory_file": "memory.json"
    }
    if not os.path.exists(path):
        with open(path, "w", encoding="utf-8") as f:
            json.dump(default_cfg, f, indent=4, ensure_ascii=False)
        return default_cfg
    with open(path, "r", encoding="utf-8") as f:
        cfg = json.load(f)
    changed = False
    for k, v in default_cfg.items():
        if k not in cfg:
            cfg[k] = v
            changed = True
    if changed:
        with open(path, "w", encoding="utf-8") as f:
            json.dump(cfg, f, indent=4, ensure_ascii=False)
    return cfg


CONFIG = None  # 延迟加载：仅 wechat_bot 独立运行模式使用（见 __main__）
# 历史缺陷：模块级 CONFIG = load_config() 在 import 时即读写 config.json——
# 任何 import（含测试、GUI）都触发磁盘 IO；GUI 模式实际用 config_store 的
# cfg（两套加载并存）。改为 __main__ 内加载，import 零副作用。


_FILE_TOKEN_RE = re.compile(
    r"[\w\u4e00-\u9fff][\w\u4e00-\u9fff\-.+()（）]*?"
    r"\.(?:docx?|xlsx?|pptx?|pdf|txt|md|html?|json|csv|zip|rar|7z|png|jpe?g|gif|mp4|mp3)",
    re.I)


def _extract_file_name_token(text):
    """从 OCR 文本拆出干净文件名 token（含常见文档扩展名，去掉大小/图标字符）。

    '部门简介+纳新宣传.docx 20.1K W' → '部门简介+纳新宣传.docx'
    文件卡片 OCR 会把文件名、大小（20.1K 带小数点）、图标字符（W/P/?）读成
    一串——整串当文件名传给 os.path.splitext 时，20.1K 的小数点会被误当
    扩展名分隔符，导致主干匹配失败。先拆出真实文件名 token。
    """
    toks = _FILE_TOKEN_RE.findall(text or "")
    return toks[0] if toks else None


# OCR 对文件名里的分隔符常漏读（真机事故：磁盘名「小漓_深海小剧场.html」
# OCR 成「小漓深海小剧场.html」——下划线在基线上，是 RapidOCR 的经典漏读
# 字符）——按名查找前把两侧的分隔符剥掉再比，容忍 OCR 丢字。
_NAME_SEP_RE = re.compile(r"[ \t\u3000_＿\-－]+")


def _norm_file_stem(stem):
    """文件名主干匹配归一化：剥掉 OCR 易漏读的分隔符（下划线/连字符/空白）。"""
    return _NAME_SEP_RE.sub("", stem or "")


def _find_file_by_display_name_impl(file_dir, display_name):
    """按消息中的显示文件名在接收目录定位对方发来的文件（模块级，单测直测）。

    微信 4.x 下载命名 '<hash>_<msgid>_m_<原名>'，目录文件名包含原名
    （分隔符归一化后包含匹配）。同名文件重复落盘时微信追加 (N) 重名后缀
    且编号单调递增——取 (N) 最大 = 最近下载的那个，平局 ctime 新者优先。

    不做 bot 发送副本排除（旧快照/成果登记方案已删）：名字锚定查找下，
    bot 的发送副本只在同名时进候选，且其编号必然小于其后用户下载产生的
    副本，(N) 最大语义天然选中用户文件——用户把 bot 发的文件发回来不再
    被误排除（真机缺陷：回传落在发送后 300s 内被「成果副本」规则拦下；
    后续又发现 OCR 漏读文件名下划线导致匹配失败，加分隔符归一化根治）。
    返回路径或 None"""
    if not display_name:
        return None
    if not file_dir or not os.path.isdir(file_dir):
        return None
    dstem = _norm_file_stem(
        re.sub(r"\(\d+\)$", "", os.path.splitext(display_name)[0]))
    if not dstem:
        return None
    best = None
    best_key = (-1, -1)  # (微信重名编号 N, ctime)：N 最大 = 最近下载
    try:
        for root, dirs, files in os.walk(file_dir):
            for fname in files:
                fstem_full = os.path.splitext(fname)[0]
                fstem = _norm_file_stem(re.sub(r"\(\d+\)$", "", fstem_full))
                if dstem not in fstem:
                    continue
                full = os.path.join(root, fname)
                try:
                    ts = os.path.getctime(full)
                except OSError:
                    continue
                m_dup = re.search(r"\((\d+)\)$", fstem_full)
                dup = int(m_dup.group(1)) if m_dup else 0
                key = (dup, ts)
                if key > best_key:
                    best_key = key
                    best = full
    except Exception as e:
        logger.error(f"[文件] 按文件名定位失败: {e}")
        return None
    if best:
        logger.info(f"[文件] 按消息文件名定位: {os.path.basename(best)}")
    return best


# 模型偶发复读进回复的开头标记：历史注入的 [time] 前缀、私聊/群聊装饰
# 「私聊 - sender」「群聊 - 群名」。剥除后再记录历史 + 发送，避免污染记忆
# 并防止后续回复继续复读。
_REPLY_PREFIX_RE = re.compile(
    r'^\s*(?:'
    r'\[\d{4}-\d{2}-\d{2}\s\d{2}:\d{2}(?::\d{2}|:xx)?\]'
    r'|\[(?:私聊|群聊)(?:\s*[-—:：]\s*[^\]]*)?\]'
    r')\s*'
)


def strip_reply_prefix(reply):
    """剥掉模型偶发复读进回复的开头标记，可叠加（如 [ts][私聊 - 王文生]），循环剥到干净。"""
    if not reply:
        return reply
    while True:
        cleaned = _REPLY_PREFIX_RE.sub("", reply, count=1)
        if cleaned == reply:
            return reply
        reply = cleaned


def _tool_call_name(tc):
    """取 tool_call 的 function.name（畸形结构返回空串）。"""
    try:
        return str((((tc or {}).get("function") or {}).get("name")) or "")
    except AttributeError:
        return ""


def _tool_arg(tc, key):
    """取工具调用 arguments JSON 里的字符串参数（解析失败返回空串）。"""
    try:
        args = json.loads(((tc or {}).get("function") or {}).get("arguments") or "{}")
    except (ValueError, TypeError):
        return ""
    return str(args.get(key) or "").strip() if isinstance(args, dict) else ""


class WeChatBot:
    def __init__(self, cfg, stop_event=None, max_connect_retries=None):
        """stop_event：微信连接重试可被外部中断（GUI 引擎停止时用，None=不中断）。
        max_connect_retries：连接失败重试上限（None=无限重试，CLI 模式保留）；
        GUI 传入有限值，超限抛异常 → 引擎进入 error 状态（可重新初始化）。"""
        self._stop_event = stop_event
        self._max_connect_retries = max_connect_retries
        self._connect_retry_interval = 10  # 微信连接失败重试间隔（秒）；测试可调小
        self.nickname = cfg["bot_nickname"]
        self.api_url = cfg["ai_api_url"]
        self.api_key = cfg["ai_api_key"]
        # 视觉端点沿用聊天端点（单模型化：call_vision_api 的 model 取 chat_model）
        self.vision_api_url = cfg.get("vision_api_url") or self.api_url
        self.vision_api_key = cfg.get("vision_api_key") or self.api_key
        self.chat_model = strip_model_prefix(cfg["chat_model"])
        self.chat_temperature = cfg.get("chat_temperature", 0.7)
        self.chat_top_p = cfg.get("chat_top_p", 0.9)
        # 单模型化后视觉调用随聊天链路：模型取 chat_model、温度取
        # chat_temperature、system 取 system_prompt（vision_prompt/vision_temp
        # 已随「图片复述」路径废弃删除）；vision_max_tokens 保留为视觉
        # max_tokens 独立上限（数据层已删键，此处 get 内置默认兜底）
        self.vision_max_tokens = cfg.get("vision_max_tokens", 10000)
        self.system_prompt = cfg.get("system_prompt", AI_DEFAULTS["system_prompt"])
        self.max_history = cfg.get("max_history", AI_DEFAULTS["max_history"])
        self.cooldown = cfg.get("cooldown", AI_DEFAULTS["cooldown"])
        self.api_retry = cfg.get("api_retry", AI_DEFAULTS["api_retry"])
        self.api_timeout = cfg.get("api_timeout", AI_DEFAULTS["api_timeout"])
        self.api_wall_budget = cfg.get("api_wall_budget", AI_DEFAULTS["api_wall_budget"])
        # 用量统计：每次 LLM 调用终态追加一行 JSONL（埋点在 _post_chat_completions）
        self.usage_store = UsageStore()
        # 微信窗口定位：None=默认右半屏；[x,y,w,h]=自定义；"off"=保持手动
        self.wechat_window_rect = cfg.get("wechat_window_rect", None)
        self.paused = cfg.get("start_paused", True)
        self.memory_file = cfg.get("memory_file", "memory.json")
        # 长记忆（v2）配置：recent 只保留 memory_keep_recent 条，溢出消息
        # append 进 memory_deep/<聊天>.jsonl 永不删除（深层记忆，recall_memory
        # 工具按需检索）；压缩 = 每攒满 memory_compress_batch 条深层消息调
        # 压缩模型提炼「重要记忆（常驻上下文）+ 关键词索引（命中时注入）」。
        # 压缩调用由后台线程（AgentBot 的 MemoryCompressor）执行，所有
        # memory_db/深层文件读写经 _memory_lock 串行化。
        self.memory_deep_enabled = bool(cfg.get("memory_deep_enabled", True))
        self.memory_compress_enabled = bool(cfg.get("memory_compress_enabled", False))
        self.memory_keep_recent = max(5, int(cfg.get("memory_keep_recent", 30)))
        self.memory_compress_batch = max(5, int(cfg.get("memory_compress_batch", 30)))
        self.memory_important_max = max(3, int(cfg.get("memory_important_max", 20)))
        self.memory_compress_model = strip_model_prefix(
            cfg.get("memory_compress_model", "") or "")
        self._deep_dir = os.path.join(
            os.path.dirname(os.path.abspath(self.memory_file)) or ".", "memory_deep")
        self._deep_count = {}   # chat -> 深层文件行数（追加时维护，启动时清点）
        self._memory_lock = threading.RLock()
        # 状态监视（set_reminder kind=condition）总开关：会产生额外 API 调用
        # （api 判定模式每轮一次小调用 + 触发回递），是否开启由用户在设置页
        # 决定（默认关）。关闭时工具分支直接友好告知，不降级普通聊天
        # （降级会让模型凭空答应没做到的提醒）。
        self.state_watch_enabled = bool(cfg.get("state_watch_enabled", False))
        # 文件处理配置
        self.file_model = strip_model_prefix(cfg.get("file_model", self.chat_model))
        self.file_temp = cfg.get("file_temp", 1.0)
        self.file_max_tokens = cfg.get("file_max_tokens", 10000)
        self.file_prompt = cfg.get("file_prompt", self.system_prompt)
        self.file_storage_path = cfg.get("file_storage_path", "")
        # 图片消息点击偏移校准（竖图点击偏位时手动校正，格式 [dx, dy]，存 config.json）
        self.image_click_offset = cfg.get("image_click_offset", [0, 0])
        # 占位消息计数（"收到任务啦，正在处理中"等）：按 chat_name 隔离，
        # 非全局计数——每个聊天独立记录，绝不跨会话复用
        self._pending_placeholders = {}

        self._model_lock = threading.RLock()

        self.memory_db = {}
        self._load_memory()
        # 节流写盘状态：高频对话下避免每条消息全量 dump memory.json
        self._memory_dirty = False
        self._last_memory_save = 0.0
        self.last_reply_time = 0
        self.wx = None
        self._connect_wx()

    def set_chat_temperature(self, value):
        with self._model_lock:
            self.chat_temperature = float(value)

    def set_chat_top_p(self, value):
        with self._model_lock:
            self.chat_top_p = float(value)

    def _load_memory(self):
        if os.path.exists(self.memory_file):
            try:
                with open(self.memory_file, "r", encoding="utf-8") as f:
                    self.memory_db = migrate_memory_data(json.load(f))
            except Exception as e:
                logger.error(f"加载记忆失败: {e}，将使用空白记忆")
                self.memory_db = {}
        # 启动一次性溢出搬运：v1 迁移来的长历史超出 recent 上限的部分
        # 全部归档进深层文件（永不删除）；同时建立深层行数计数
        cap = self._recent_cap()
        for chat, st in self.memory_db.items():
            if len(st["recent"]) > cap:
                for msg in st["recent"][:-cap]:
                    self._append_deep(chat, msg)
                st["recent"] = st["recent"][-cap:]
            self._deep_count[chat] = self._count_deep(chat)

    def _recent_cap(self):
        """近期记忆保留条数：memory_keep_recent 与 max_history 取小
        （角色卡 max_history 若被用户调小仍然生效）。"""
        return max(5, min(int(getattr(self, "memory_keep_recent", 30)),
                          int(getattr(self, "max_history", 1000))))

    def _save_memory(self):
        try:
            with open(self.memory_file, "w", encoding="utf-8") as f:
                json.dump(self.memory_db, f, ensure_ascii=False, indent=2)
        except Exception as e:
            logger.error(f"保存记忆失败: {e}")

    def _connect_wx(self):
        """连接微信（后端抽象层）。失败按 _connect_retry_interval 间隔重试：
        - stop_event 已设置 → 立即抛异常（引擎停止/取消初始化）
        - max_connect_retries 达到 → 抛异常（GUI 初始化超时 → error 状态）
        睡眠分片（0.5s）保证 stop 响应延迟 ≤0.5s。"""
        attempts = 0
        while True:
            if self._stop_event is not None and self._stop_event.is_set():
                raise RuntimeError("微信连接已取消")
            try:
                self.wx = create_backend("auto")
                self._init_position_wechat()
                logger.info(f"✅ 微信连接成功（后端: {self.wx.name}）")
                return
            except Exception as e:
                attempts += 1
                if self._max_connect_retries is not None \
                        and attempts >= self._max_connect_retries:
                    raise RuntimeError(
                        f"微信连接失败（已重试 {attempts} 次，超过上限 "
                        f"{self._max_connect_retries}）：{e}") from e
                logger.error(f"微信连接失败：{e}，{self._connect_retry_interval}秒后重试...")
                remain = self._connect_retry_interval
                while remain > 1e-9:
                    if self._stop_event is not None and self._stop_event.is_set():
                        raise RuntimeError("微信连接已取消")
                    time.sleep(min(0.5, remain))
                    remain -= 0.5

    def _init_position_wechat(self):
        """初始化微信窗口定位（连接成功时执行一次，用户定案）。

        视觉坐标依赖窗口位置/尺寸稳定——初始化即摆到标准位并告知用户
        勿动，替代旧「不移动窗口」约定。配置 wechat_window_rect：
        - None（默认）→ 主屏右半边（README 系统要求）
        - [x, y, w, h] → 自定义矩形（物理像素）
        - "off" → 完全保持手动（旧行为）
        运行期不再核对矩形（用户否决每轮像素核对的成本），仅保留最小化
        哨兵自动恢复（IsIconic 一次系统调用，微秒级）。
        """
        rect = getattr(self, "wechat_window_rect", None)
        if rect == "off":
            logger.info("[定位] wechat_window_rect=off，保持手动窗口位置")
            return
        hwnd = find_window_by_title("微信")
        if not hwnd:
            logger.warning("[定位] 未找到微信窗口，跳过定位")
            return
        ensure_window_visible(hwnd)  # 最小化先拉起，定位才有意义
        try:
            if ctypes.windll.user32.IsZoomed(hwnd):
                # 最大化先还原再定位：SetWindowPos 对最大化窗口行为不可靠
                # （可能只改尺寸不改状态，视觉错乱）
                ctypes.windll.user32.ShowWindow(hwnd, 9)  # SW_RESTORE
                time.sleep(0.3)
        except Exception:
            pass
        if isinstance(rect, (list, tuple)) and len(rect) == 4:
            try:
                x, y, w, h = (int(v) for v in rect)
                source = "配置"
            except (TypeError, ValueError):
                x, y, w, h = default_right_half_rect()
                source = "配置非法，回退默认右半屏"
        else:
            x, y, w, h = default_right_half_rect()
            source = "默认右半屏"
        if position_window_visible(hwnd, x, y, w, h):
            # 可见内容语义：自动外扩 DWM 不可见边框，可见部分精确贴合目标
            # （与手动拖窗口到打满的系统行为一致；直接摆窗口矩形会两侧
            # 各缩进 ~10px——用户实测「右侧有缝隙」）
            logger.info(f"[定位] 微信窗口已定位（{source}）：可见区 ({x},{y}) {w}x{h}"
                        "——请勿最小化或调整窗口大小，否则影响消息识别")
        else:
            logger.warning("[定位] 微信窗口定位失败，保持当前位置")

    def _chat_state(self, chat_id):
        """取（或创建）聊天的 v2 记忆状态 {recent, important, index, indexed}。
        调用方必须已持 _memory_lock。"""
        st = self.memory_db.get(chat_id)
        if st is None:
            st = {"recent": [], "important": [], "index": [], "indexed": 0}
            self.memory_db[chat_id] = st
        return st

    def _get_history(self, chat_id):
        """近期记忆（注入上下文的 recent 窗口）。深层历史不在此列——由
        recall_memory 工具按需检索。"""
        chat_id = _memory_key(chat_id)
        with self._memory_lock:
            return self._chat_state(chat_id)["recent"]

    def _add_history(self, chat_id, role, content):
        chat_id = _memory_key(chat_id)
        with self._memory_lock:
            st = self._chat_state(chat_id)
            st["recent"].append({
                "role": role,
                "content": content,
                "time": time.strftime("%Y-%m-%d %H:%M:%S")
            })
            cap = self._recent_cap()
            while len(st["recent"]) > cap:
                # 溢出归档：深层记忆启用时写入 memory_deep/ 永久保存，
                # 未启用则与旧行为一致（超出窗口即丢弃）
                self._append_deep(chat_id, st["recent"].pop(0))
            self._schedule_save_memory()

    # ---------- 深层记忆（memory_deep/<chat>.jsonl，append-only 永不删除） ----------

    def _deep_path(self, chat_id):
        """深层记忆文件路径：聊天名 percent-encode（文件名安全且可逆）。"""
        from urllib.parse import quote
        return os.path.join(self._deep_dir, quote(chat_id, safe="") + ".jsonl")

    def _append_deep(self, chat_id, msg):
        """一条消息溢出 recent 时归档进深层文件（启用深层记忆才写）。"""
        if not getattr(self, "memory_deep_enabled", False):
            return
        try:
            os.makedirs(self._deep_dir, exist_ok=True)
            with open(self._deep_path(chat_id), "a", encoding="utf-8") as f:
                f.write(json.dumps(msg, ensure_ascii=False) + "\n")
            self._deep_count[chat_id] = self._deep_count.get(chat_id, 0) + 1
        except Exception as e:
            logger.error(f"[记忆] 深层写入失败: {e}")

    def _count_deep(self, chat_id):
        """清点深层文件行数（启动时建立计数基线；坏行跳过不计数）。"""
        n = 0
        try:
            with open(self._deep_path(chat_id), "r", encoding="utf-8") as f:
                for line in f:
                    line = line.strip()
                    if line and line.startswith("{"):
                        n += 1
        except OSError:
            return 0
        return n

    def _iter_deep(self, chat_id):
        """按序读取深层文件全部消息（坏行跳过）。"""
        try:
            with open(self._deep_path(chat_id), "r", encoding="utf-8") as f:
                for line in f:
                    line = line.strip()
                    if not line:
                        continue
                    try:
                        item = json.loads(line)
                    except ValueError:
                        continue
                    if isinstance(item, dict):
                        yield item
        except OSError:
            return

    def _read_deep_range(self, chat_id, start, count):
        """读取深层文件 [start, start+count) 区间的消息（压缩线程取待压缩段）。
        返回 (消息列表, 实际区间终点)：文件行数少于预期（文件被外部改动）
        时按实际返回。"""
        out = []
        for i, msg in enumerate(self._iter_deep(chat_id)):
            if i < start:
                continue
            out.append(msg)
            if len(out) >= count:
                break
        return out, start + len(out)

    def _clear_deep(self, chat_id=None):
        """删除深层记忆文件（清空记忆时联动；chat_id=None 清全部）。"""
        try:
            if chat_id:
                paths = [self._deep_path(chat_id)]
            else:
                paths = [os.path.join(self._deep_dir, n)
                         for n in os.listdir(self._deep_dir)
                         if n.endswith(".jsonl")] if os.path.isdir(self._deep_dir) else []
            for p in paths:
                if os.path.isfile(p):
                    os.remove(p)
        except OSError as e:
            logger.error(f"[记忆] 深层清理失败: {e}")
        if chat_id:
            self._deep_count.pop(chat_id, None)
        else:
            self._deep_count.clear()

    # ---------- recall_memory 工具：深层 + 近期记忆关键词检索 ----------

    def _recall_memory(self, chat_id, query):
        """在本聊天的全部记忆（深层存档 + 近期窗口）中按关键词检索。
        返回给模型的可读文本（带时间戳的命中消息列表）。"""
        if not getattr(self, "memory_deep_enabled", False):
            return "深层记忆未启用，无法检索历史"
        chat_id = _memory_key(chat_id)
        terms = [t for t in re.split(r"[\s，。！？、：；\"'（）()\[\]【】]+",
                                     str(query or "").strip()) if t]
        if not terms:
            return "检索词为空，请给出人名/事件/物品等关键词"
        terms_l = [t.lower() for t in terms]

        def _hit(msg):
            c = str(msg.get("content") or "").lower()
            return any(t in c for t in terms_l)

        with self._memory_lock:
            hits = [m for m in self._iter_deep(chat_id) if _hit(m)]
            hits += [m for m in self._chat_state(chat_id)["recent"] if _hit(m)]
        if not hits:
            return "没有找到相关记忆"
        tail = hits[-20:]  # 最多给模型 20 条（时间最晚的优先）
        lines = []
        for m in tail:
            role = "用户" if m.get("role") == "user" else "小漓"
            lines.append(f"[{m.get('time', '')}] {role}: "
                         f"{str(m.get('content') or '')[:150]}")
        head = (f"共命中 {len(hits)} 条，显示最近的 {len(tail)} 条：\n"
                if len(hits) > len(tail) else "")
        return head + "\n".join(lines)

    # ---------- 压缩产出：重要记忆（常驻注入）+ 关键词索引（命中注入） ----------

    def _important_block(self, chat_id):
        """重要记忆 system 块（无内容返回 None）。per-chat 独立 system——
        不得并进人设消息（人设是跨聊天共享的缓存前缀）。"""
        if not getattr(self, "memory_compress_enabled", False):
            return None
        st = self.memory_db.get(_memory_key(chat_id)) or {}
        items = [str(x.get("content") or "").strip()
                 for x in (st.get("important") or []) if isinstance(x, dict)]
        items = [x for x in items if x]
        if not items:
            return None
        return ("以下是你在与该联系人长期相处中沉淀的重要记忆，回复时遵循：\n"
                + "\n".join(f"{i}. {c}" for i, c in enumerate(items, 1)))

    def _match_related_memory(self, chat_id, user_text):
        """关键词索引匹配：用户消息命中索引关键词时返回相关记忆注入块。"""
        if not getattr(self, "memory_compress_enabled", False):
            return None
        text = str(user_text or "")
        if not text.strip():
            return None
        st = self.memory_db.get(_memory_key(chat_id)) or {}
        hits = []
        for e in (st.get("index") or []):
            if not isinstance(e, dict):
                continue
            kws = [str(k) for k in (e.get("kw") or []) if str(k).strip()]
            mem = str(e.get("mem") or "").strip()
            if mem and any(k in text for k in kws):
                hits.append(mem)
                if len(hits) >= 3:
                    break
        if not hits:
            return None
        return ("[相关记忆]（历史对话中与本次消息相关的内容，供参考）\n"
                + "\n".join(f"- {h}" for h in hits))

    def memory_commit_compression(self, chat_id, consumed, important_new, index_new):
        """压缩线程写回：推进 indexed 边界 + 合并重要记忆/索引条目（锁内）。

        consumed：本次已压缩到的深层行数（绝对值）；important_new：
        [{"content": ...}]；index_new：[{"kw": [...], "mem": ...}]。
        条目超上限时裁掉最旧的（重要记忆上限 memory_important_max，
        索引上限 300——索引只增会让注入匹配越来越慢且陈旧）。"""
        chat_id = _memory_key(chat_id)
        with self._memory_lock:
            st = self._chat_state(chat_id)
            st["indexed"] = max(int(st.get("indexed") or 0), int(consumed))
            for item in important_new:
                if isinstance(item, dict) and str(item.get("content") or "").strip():
                    item = dict(item, time=time.strftime("%Y-%m-%d %H:%M:%S"))
                    st["important"].append(item)
            limit = int(getattr(self, "memory_important_max", 20))
            st["important"] = st["important"][-limit:]
            for item in index_new:
                if isinstance(item, dict) \
                        and [k for k in (item.get("kw") or []) if str(k).strip()] \
                        and str(item.get("mem") or "").strip():
                    st["index"].append(item)
            st["index"] = st["index"][-300:]
            self._schedule_save_memory()

    # ---------- 记忆管理页的数据通道（运行中经 engine.bot 调用；全部持锁） ----------

    def memory_overview(self):
        """全部聊天的计数快照（UI 线程读，锁内构建）：{chat: {recent, deep,
        important, index}}。deep 取 _deep_count（深层文件有效行数）。"""
        with self._memory_lock:
            out = {}
            for chat, st in self.memory_db.items():
                out[chat] = {
                    "recent": len(st.get("recent") or []),
                    "deep": int(self._deep_count.get(chat, 0)),
                    "important": len(st.get("important") or []),
                    "index": len(st.get("index") or []),
                }
            return out

    def memory_detail(self, chat_id, deep_offset=0, deep_limit=200,
                      deep_query=None):
        """单聊天详情快照（锁内拷贝，UI 展示用）。

        deep_query 非空：深层全量过滤，deep 截 200 条 + deep_matched 总数；
        否则深层返回第 deep_offset 页（deep_limit 条），deep_matched=None。"""
        chat_id = _memory_key(chat_id)
        with self._memory_lock:
            st = self.memory_db.get(chat_id) or {}
            deep_total = int(self._deep_count.get(chat_id, 0))
            deep_matched = None
            if deep_query:
                q = str(deep_query).lower()
                hits = [dict(m) for m in self._iter_deep(chat_id)
                        if q in str(m.get("content") or "").lower()]
                deep_matched = len(hits)
                deep = hits[:200]
            else:
                deep = []
                for i, m in enumerate(self._iter_deep(chat_id)):
                    if i >= deep_offset and len(deep) < deep_limit:
                        deep.append(dict(m))
            return {
                "recent": [dict(m) for m in (st.get("recent") or [])],
                "important": [dict(m) for m in (st.get("important") or [])],
                "index": [dict(m) for m in (st.get("index") or [])],
                "deep_total": deep_total,
                "deep": deep,
                "deep_matched": deep_matched,
            }

    def delete_important(self, chat_id, idx):
        """删除第 idx（1 基）条重要记忆。返回是否删除。"""
        chat_id = _memory_key(chat_id)
        with self._memory_lock:
            st = self._chat_state(chat_id)
            if 1 <= idx <= len(st["important"]):
                st["important"].pop(idx - 1)
                self._schedule_save_memory()
                return True
            return False

    def delete_index_entry(self, chat_id, idx):
        """删除第 idx（1 基）条关键词索引。返回是否删除。"""
        chat_id = _memory_key(chat_id)
        with self._memory_lock:
            st = self._chat_state(chat_id)
            if 1 <= idx <= len(st["index"]):
                st["index"].pop(idx - 1)
                self._schedule_save_memory()
                return True
            return False

    def delete_deep_message(self, chat_id, line_no):
        """删除深层存档第 line_no（1 基，时间正序）条消息。

        深层文件原子重写；行号落在压缩边界 indexed 之前时 indexed 同步
        -1（边界按行数推进，少一行必须回退，否则下轮压缩错位跳过一条）。
        返回是否删除。"""
        chat_id = _memory_key(chat_id)
        with self._memory_lock:
            st = self._chat_state(chat_id)
            removed, new_count = deep_delete_line(self._deep_path(chat_id),
                                                  line_no)
            if removed is None:
                return False
            self._deep_count[chat_id] = new_count
            if line_no - 1 < int(st.get("indexed") or 0):
                st["indexed"] = int(st["indexed"]) - 1
            self._schedule_save_memory()
            return True

    def _schedule_save_memory(self):
        """节流写盘：距上次写盘 ≥1s 立即写，否则只标记脏。
        兜底时钟在 _flush_memory_if_due（引擎每 0.5s 轮询时检查到期）——
        稀疏对话下脏数据最多 ~1.5s 落盘，memory.json 查看不再滞后到
        下一条消息（历史缺陷：只标记脏无定时器，最后一条回复要等下一条
        消息才落盘）。"""
        self._memory_dirty = True
        now = time.time()
        if now - getattr(self, "_last_memory_save", 0.0) >= 1.0:
            self._flush_memory()

    def _flush_memory_if_due(self):
        """节流兜底：脏数据超过 1s 未落盘就强制写（引擎轮询每 0.5s 调用）。

        引擎主循环顶部调用本检查即成为兜底时钟——不引入线程/锁，仍全部
        在引擎线程上执行。"""
        if getattr(self, "_memory_dirty", False) \
                and time.time() - getattr(self, "_last_memory_save", 0.0) >= 1.0:
            self._flush_memory()

    def _flush_memory(self):
        """有脏数据则写盘。程序退出/引擎停止前调用，保证最近消息不丢。"""
        if getattr(self, "_memory_dirty", False):
            self._memory_dirty = False
            self._last_memory_save = time.time()
            self._save_memory()

    def clear_history(self, chat_id=None):
        with self._memory_lock:
            if chat_id:
                key = _memory_key(chat_id)
                self.memory_db.pop(key, None)
                self._clear_deep(key)
                logger.info(f"已清空聊天 {chat_id} 的历史（含深层记忆）")
            else:
                self.memory_db.clear()
                self._clear_deep()
                logger.info("已清空全部对话历史（含深层记忆）")
            self._save_memory()

    def delete_messages(self, chat_id, indices):
        chat_id = _memory_key(chat_id)
        with self._memory_lock:
            st = self.memory_db.get(chat_id)
            if st is None:
                logger.warning(f"❌ 聊天 {chat_id} 不存在于记忆中")
                return False
            hist = st["recent"]
            total = len(hist)
            to_delete = []
            for idx in indices:
                if 1 <= idx <= total:
                    to_delete.append(idx - 1)
                else:
                    logger.warning(f"序号 {idx} 超出范围（1-{total}），已忽略")
            if not to_delete:
                return False
            to_delete = sorted(set(to_delete), reverse=True)
            deleted_msgs = []
            for i in to_delete:
                deleted_msgs.append(hist.pop(i))
            self._save_memory()
        logger.info(f"已从聊天 {chat_id} 中删除 {len(deleted_msgs)} 条消息")
        for msg in deleted_msgs:
            role = "用户" if msg["role"] == "user" else "小漓"
            ts = msg.get("time", "未知时间")
            logger.info(f"  删除: [{ts}] {role}: {msg['content'][:50]}...")
        return True

    def _post_chat_completions(self, url, headers, payload, timeout, label="api", meta=None):
        """OpenAI 兼容 chat/completions 统一调用入口（chat / vision 两链路共用）。

        重试策略：
        - 429/5xx/网络异常/坏 JSON → 指数退避重试（1s 起步、封顶 8s）；
          服务器带 Retry-After 时尊重之，超过 RETRY_AFTER_GIVEUP 直接放弃
        - 其余 4xx（鉴权/参数错误）不重试——重试无意义，只会拖慢用户感知
        - 墙钟预算 api_wall_budget（默认 45s）：重试总时长封顶。历史缺陷：
          每次超时 60s × api_retry 3 次 = 用户干等 3 分钟才收到报错
        成功返回解析后的 dict；最终失败抛 ApiCallError。绝不返回
        「API 错误: xxx」这类会原样发给好友的字符串（历史行为，已废）。
        meta：dict（kind/model/messages），透传给用量统计埋点。
        """
        wall_budget = float(getattr(self, "api_wall_budget", API_WALL_BUDGET_DEFAULT))
        deadline = time.monotonic() + wall_budget
        attempts = int(getattr(self, "api_retry", 2)) + 1
        latency_ms = 0.0
        last_desc = "未发起请求"
        status = None
        for attempt in range(attempts):
            if attempt > 0:
                remain = deadline - time.monotonic()
                if remain <= 0:
                    last_desc = f"墙钟预算耗尽（{wall_budget:.0f}s），停止重试"
                    break
            started = time.monotonic()
            retryable = False
            retry_after = None
            data = None
            try:
                resp = requests.post(url, headers=headers, json=payload, timeout=timeout)
            except requests.exceptions.RequestException as e:
                latency_ms += (time.monotonic() - started) * 1000.0
                last_desc = f"网络异常 {type(e).__name__}: {e}"
                logger.warning(f"[{label}] {last_desc}（第 {attempt + 1} 次）")
                retryable = True
            else:
                status = resp.status_code
                latency_ms += (time.monotonic() - started) * 1000.0
                if status == 200:
                    try:
                        data = resp.json()
                    except ValueError as e:
                        last_desc = f"响应 JSON 解析失败: {e}"
                        logger.warning(f"[{label}] {last_desc}")
                        retryable = True
                    else:
                        self._finish_usage(meta, ok=True, status=200,
                                           latency_ms=latency_ms, data=data)
                        return data
                else:
                    last_desc = f"HTTP {status}: {(getattr(resp, 'text', '') or '')[:200]}"
                    logger.warning(f"[{label}] {last_desc}")
                    if status == 429 or 500 <= status < 600:
                        retryable = True
                        resp_headers = getattr(resp, "headers", None) or {}
                        for key in resp_headers:
                            if str(key).lower() == "retry-after":
                                try:
                                    retry_after = float(resp_headers[key])
                                except (TypeError, ValueError):
                                    retry_after = None
                                break
                    else:
                        self._finish_usage(meta, ok=False, status=status,
                                           latency_ms=latency_ms)
                        raise ApiCallError(last_desc, status=status)
            if retryable and retry_after is not None and retry_after > RETRY_AFTER_GIVEUP:
                self._finish_usage(meta, ok=False, status=status, latency_ms=latency_ms)
                raise ApiCallError(
                    f"HTTP {status}（Retry-After {retry_after:.0f}s 超过放弃阈值）",
                    status=status)
            if attempt + 1 < attempts:
                delay = (retry_after if retry_after is not None
                         else min(BACKOFF_CAP, BACKOFF_BASE * (2.0 ** attempt)))
                delay = min(delay, max(0.0, deadline - time.monotonic()))
                if delay > 0:
                    time.sleep(delay)
        self._finish_usage(meta, ok=False, status=status, latency_ms=latency_ms)
        raise ApiCallError(f"{label} 调用失败: {last_desc}", status=status)

    def _finish_usage(self, meta, ok, status, latency_ms, data=None):
        """用量统计埋点（_post_chat_completions 终态调用一次）。

        usage_store 缺失（测试桩 bot / 未初始化）时静默跳过。API 响应缺
        usage 字段时用 estimate_tokens 兜底估算。缓存字段（命中/未命中/
        推理 tokens）从响应 usage 透传：DeepSeek 用顶层 prompt_cache_hit/
        miss_tokens，OpenAI 系用 prompt_tokens_details.cached_tokens 与
        completion_tokens_details.reasoning_tokens——v2.5.0 的消息布局缓存
        优化效果就靠这几列验证。src 标记数据来源：api=响应实测，est=本地
        估算（可信度分开，不混算）。"""
        store = getattr(self, "usage_store", None)
        if store is None:
            return
        try:
            meta = dict(meta or {})
            usage = (data or {}).get("usage") or {}
            prompt_t = usage.get("prompt_tokens")
            completion_t = usage.get("completion_tokens")
            if prompt_t is None and meta.get("messages"):
                prompt_t = estimate_tokens("".join(
                    str(m.get("content") or "") for m in meta["messages"]))
            if completion_t is None and data:
                try:
                    completion_t = estimate_tokens(
                        data["choices"][0]["message"]["content"] or "")
                except (KeyError, IndexError, TypeError):
                    pass

            def _int(v):
                try:
                    return int(v)
                except (TypeError, ValueError):
                    return 0

            p_details = usage.get("prompt_tokens_details") or {}
            c_details = usage.get("completion_tokens_details") or {}
            cache_hit = usage.get("prompt_cache_hit_tokens")
            if cache_hit is None:
                cache_hit = p_details.get("cached_tokens")
            cache_hit = _int(cache_hit)
            cache_miss = usage.get("prompt_cache_miss_tokens")
            if cache_miss is None and prompt_t is not None:
                cache_miss = _int(prompt_t) - cache_hit
            else:
                cache_miss = _int(cache_miss)
            store.record(kind=meta.get("kind"), model=meta.get("model"),
                         prompt_tokens=prompt_t, completion_tokens=completion_t,
                         ok=ok, status=status, latency_ms=latency_ms,
                         cache_hit=cache_hit, cache_miss=cache_miss,
                         reasoning=_int(c_details.get("reasoning_tokens")),
                         total_tokens=usage.get("total_tokens"),
                         src="api" if (usage.get("prompt_tokens") is not None
                                       or usage.get("completion_tokens") is not None)
                         else "est")
        except Exception as e:
            logger.debug(f"[用量] 记录失败: {e}")

    def call_vision_api(self, content, chat_id=None, related_memory=None):
        """单调用视觉识别（OpenAI 兼容 / chat.completions）。

        content：块列表 list[dict]，格式
          [{"type": "text", "text": ...},
           {"type": "image_url", "image_url": {"url": "data:image/..."}}]
        图片块可选——无图时只含 text 块（调用方构造，本方法原样透传进 user
        消息；DeepSeek vision 限制：图片只能出现在 user 消息，system/assistant
        带图返回 400）。消息布局（缓存友好：稳定前缀在前、每轮变化区在尾）：
          [system 人设] → [system 重要记忆] → [历史(带[ts])] →
          [system 相关记忆] → [system 当前时间] → [user content]
        人设（self.system_prompt）前置为 system 纯文本消息（空人设则不插入人设
        system 消息）；重要记忆/相关记忆来自长记忆压缩产出（per-chat，未启用
        或无内容时缺省）；「当前时间」system 紧贴当前消息——它每秒变化，绝不能
        插在历史之前打断缓存前缀。

        chat_id 可选（默认 None）：非空时注入重要记忆块与 _get_history(chat_id)
        近期历史（语义逐字对齐 call_chat_ai：有 time 字段带 [ts] 前缀，否则
        原文；不重排）。为空时仍注入当前时间 system（图片/文件描述路径自动
        受益；persona 为空时本条保证 messages 至少一条 system）。

        payload 声明 dispatch_task / set_reminder / recall_memory / web_search /
        web_fetch 工具（tool_choice=auto）。web_search/web_fetch/recall_memory
        是查询型工具循环：模型调用后本方法执行（搜索/抓正文/记忆检索），把
        结果作为 tool 消息回填、继续补全，直到给出最终答复。整个调用最多
        VISION_TOOL_ROUNDS 次补全，循环耗尽仍未收敛返回 None（调用方降级）。
        dispatch_task/set_reminder 调用优先原样返回（既有单次语义不变，查询
        往返不延迟任务投递）。工具往返只存在于本次调用的 messages，不写入
        对话历史。

        返回结构化结果（供上层按 dict 处理）：
        - message.tool_calls 含 dispatch_task/set_reminder →
          {'kind': 'tool_call', 'name': ..., 'arguments': ...}（arguments 为原始 JSON 字符串）
        - 仅 message.content（含经工具循环后的最终答复）→
          {'kind': 'text', 'content': ...}
        - 非 200 / 无 choices / content 空白 / 循环耗尽 → None
        """
        headers = {"Authorization": f"Bearer {self.vision_api_key}", "Content-Type": "application/json"}
        # 人设由 system 纯文本消息承载（绝不放图片——DeepSeek 限制图片
        # 只能进 user 消息）；persona 为空时不插入空 system 消息。
        persona = (getattr(self, "system_prompt", "") or "").strip()
        messages = [{"role": "system", "content": persona}] if persona else []
        important = self._important_block(chat_id) if chat_id else None
        if important:
            messages.append({"role": "system", "content": important})
        if chat_id:
            # 历史注入：语义逐字对齐 call_chat_ai（system 之后、user 之前；
            # 有 time 字段带 [ts] 前缀，否则原文；不重排——_get_history 返回
            # 列表本身已按时间有序）
            for h in self._get_history(chat_id):
                if "time" in h:
                    ts = h["time"]
                    msg_content = f"[{ts}] {h['content']}"
                else:
                    msg_content = h['content']
                messages.append({"role": h["role"], "content": msg_content})
        if related_memory:
            messages.append({"role": "system", "content": related_memory})
        # 当前时间 system：无条件注入且紧贴当前消息（历史之后——它每秒变化，
        # 插在历史前会打断缓存前缀）；persona 为空时本条保证 messages 至少
        # 一条 system，消除空 messages 隐患。
        current_time = time.strftime("%Y-%m-%d %H:%M:%S")
        messages.append({"role": "system", "content": f"当前时间：{current_time}"})
        messages.append({"role": "user", "content": content})
        with self._model_lock:
            # 单模型化：视觉 model 取 chat_model（__init__ 已 strip 前缀），
            # 空则兜底 vision-exp 纯名（防空 model / 带前缀兜底 → API 400）；
            # 温度随聊天温度（vision_temp 已删）
            model = self.chat_model or VISION_MODEL_DEFAULT
            temp = self.chat_temperature
        # 上下文预算裁剪：超长历史/文件全文会撑爆模型上下文上限
        # （实测请求 272 万 token → API 400 "maximum context length"）。
        # 逐字对齐 call_chat_ai：从最旧历史开始丢弃，保证单次请求不超模型上下文。
        messages = fit_messages_in_budget(
            messages, budget=getattr(self, "max_context_tokens", 100000))
        tools = [{
            "type": "function",
            "function": {
                "name": "dispatch_task",
                "description": "判断用户消息是否为任务，是则投递天枢处理",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "task": {"type": "string", "description": "任务描述"},
                    },
                    "required": ["task"],
                },
            },
        }]
        if getattr(self, "reminders", None) is not None:
            # 统一触发器工具（kind=time 定时 / kind=condition 状态监视）：
            # 模型依据 system 的「当前时间」消息把相对表达换算成绝对时间。
            # 无 content 参数——到点/达成只回递事件，回复由模型结合创建时
            # 的历史对话自行生成（用户定案：不做预写文案的闹钟）。
            # 状态监视纪律（真机实验定案，见 ConditionWatcher）：创建前必须
            # 用 web_search 挑页面并 web_fetch 看过实际内容，轮询期只访问
            # 固定 URL、不再调用搜索引擎；页面混合当前与未来时段（如天气预报
            # 页「今天暴雨…周六晴」）时必须给 scope 切片，否则本地判定会被
            # 未来预报误触发。
            tools.append({
                "type": "function",
                "function": {
                    "name": "set_reminder",
                    "description": "创建触发器，两种模式二选一（到点/达成时系统会把"
                                   "触发事件回递给你，由你结合创建该触发器时的对话"
                                   "历史按人设自然回复，因此不需要预先写提醒文案）：\n"
                                   "① kind=time 定时：到指定时间触发"
                                   "（如「明天下午3点提醒我查成绩」「每天早上8点叫我起床」）。\n"
                                   "② kind=condition 状态监视：轮询监控一个固定网页，"
                                   "条件达成或到截止时间时触发"
                                   "（如「如果一会儿福州大学旗山校区下雨了提醒我去拿快递」"
                                   "「盯着一个商品页面，降价了告诉我」）。\n"
                                   "状态监视创建规则：先用 web_search 搜索并挑一个内容"
                                   "稳定的页面，用 web_fetch 看过实际内容后再创建"
                                   "（关键词要按页面实际措辞给全同义形式）；轮询期只"
                                   "访问该固定 URL、不再调用搜索引擎；若页面同时包含"
                                   "当前与未来时段的信息（如天气预报页），必须给"
                                   "scope_start/scope_end 圈出与条件相关的当前时段切片"
                                   "（如「今天」到「明天」），绝不能整页匹配；条件无法"
                                   "用关键词判断（需要理解语义或比较数值）时用 judge=api。",
                    "parameters": {
                        "type": "object",
                        "properties": {
                            "kind": {"type": "string", "enum": ["time", "condition"],
                                     "description": "触发器类型：time=定时，"
                                                    "condition=状态监视"},
                            "time": {"type": "string",
                                     "description": "kind=time 必填。触发时间，格式"
                                                    " YYYY-MM-DD HH:MM（24 小时制；"
                                                    "依据当前时间换算相对表达）"},
                            "repeat": {"type": "string", "enum": ["once", "daily", "weekly"],
                                       "description": "kind=time：重复规则，默认 once"},
                            "condition": {"type": "string",
                                          "description": "kind=condition 必填。自然语言"
                                                         "条件原文（判定与到期回递的依据）"},
                            "url": {"type": "string",
                                    "description": "kind=condition 必填。要固定轮询监控"
                                                   "的网页 URL（创建时已用 web_fetch 验看过）"},
                            "judge": {"type": "string", "enum": ["local", "api"],
                                      "description": "kind=condition：达成判定方式，默认"
                                                     " local。local=本地关键词匹配（零 API"
                                                     " 成本，适合文字标记型条件）；api=每次"
                                                     "抓取后调模型判定（适合需理解语义或比较"
                                                     "数值的条件）"},
                            "match_type": {"type": "string", "enum": ["present", "absent"],
                                           "description": "kind=condition 且 judge=local："
                                                          "present=切片内出现任一关键词即达成"
                                                          "（如「有货」）；absent=切片内关键词"
                                                          "全部消失即达成（如「等雨停」盯「雨」"
                                                          "消失）。默认 present"},
                            "met_keywords": {"type": "array",
                                             "items": {"type": "string"},
                                             "description": "kind=condition 且 judge=local"
                                                            " 必填：判定关键词数组（如"
                                                            " [\"雨\",\"降雨\",\"雷阵雨\"]），"
                                                            "按页面实际措辞给全同义形式"},
                            "scope_start": {"type": "string",
                                            "description": "kind=condition 可选：切片起始"
                                                           "标记（页面文本原文，如「今天」）"},
                            "scope_end": {"type": "string",
                                          "description": "kind=condition 可选：切片结束标记"
                                                         "（如「明天」），只匹配两标记之间，"
                                                         "隔离未来时段"},
                            "interval_seconds": {"type": "integer",
                                                 "description": "kind=condition：轮询间隔秒，"
                                                                "默认 60，最小 10"},
                            "expire_at": {"type": "string",
                                          "description": "kind=condition 可选：截止时间"
                                                         " YYYY-MM-DD HH:MM，不填默认 7 天后"},
                        },
                        "required": ["kind"],
                    },
                },
            })
        # 联网搜索工具（零配置：百度/必应/搜狗并发合并，见 xiaoli_app/web_search）。
        # query 描述写明真机校准的措辞规则：多个主体用空格分开（百度/搜狗对
        # 「机构 人名」式查询精准命中，cn.bing 反而会分词失败）；并给出结果
        # 无关时的换措辞重试策略——模型拿到垃圾结果只会硬答或放弃是此前
        # 搜索「不精确」的主因之一。
        tools.append({
            "type": "function",
            "function": {
                "name": "web_search",
                "description": "联网搜索（搜狗/必应/百度多引擎并发）。需要实时信息"
                               "（天气/新闻/价格/赛事结果等）或拿不准的事实时调用。"
                               "query 一次只查一个主题，用简短中文短语，多个主体用"
                               "空格分开（如「福州天气」「福州大学 唐勇」），不要堆"
                               "修饰词。若返回结果与查询主题明显无关（如搜人名得到"
                               "无关实体），换措辞再搜一次：调整词序、增删限定词"
                               "（如加「教授」「简介」）或改用更具体的表述",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "query": {"type": "string",
                                  "description": "简短中文短语，一次一个主题，"
                                                 "多个主体用空格分开"
                                                 "（如「福州天气」「福州大学 唐勇」）"},
                    },
                    "required": ["query"],
                },
            },
        })
        # 网页正文抓取（与 web_search 成对）：搜索摘要只有站点介绍，实时数据
        # 在网页正文里——模型从搜索结果挑来源后抓正文自己读
        tools.append({
            "type": "function",
            "function": {
                "name": "web_fetch",
                "description": "抓取网页正文全文。搜索结果里挑最相关的来源"
                               "（如天气网页面）后调用，读取其中具体的数据"
                               "（气温/天气/比分等）",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "url": {"type": "string", "description": "要抓取的网页 URL"},
                    },
                    "required": ["url"],
                },
            },
        })
        # 深层记忆检索工具：仅启用深层记忆且有聊天上下文时声明——模型在
        # 「用户提到过去的事但自己记不清」时主动调用，避免凭空编造
        if chat_id and getattr(self, "memory_deep_enabled", False):
            tools.append({
                "type": "function",
                "function": {
                    "name": "recall_memory",
                    "description": "在本聊天的全部历史记忆中检索（包括很久以前的"
                                   "对话存档）。当用户提到过去聊过的事、之前的约定"
                                   "或承诺、或问你记不记得某件事而你一时想不起来时"
                                   "调用，不要凭空编造",
                    "parameters": {
                        "type": "object",
                        "properties": {
                            "query": {"type": "string",
                                      "description": "检索关键词，如人名/事件/物品名"
                                                     "（可给多个，空格分开）"},
                        },
                        "required": ["query"],
                    },
                },
            })
        payload = {
            "model": model,
            "messages": messages,
            "max_tokens": self.vision_max_tokens,
            "temperature": temp,
            "tools": tools,
            "tool_choice": "auto",
        }
        try:
            for _round in range(VISION_TOOL_ROUNDS):
                data = self._post_chat_completions(
                    self.vision_api_url, headers, payload, 60, label="vision",
                    meta={"kind": "vision", "model": model, "messages": messages})
                choices = data.get("choices", [])
                if not choices:
                    logger.warning("视觉模型返回无 choices")
                    return None
                message = choices[0].get("message", {}) or {}
                tool_calls = message.get("tool_calls") or []
                if not tool_calls:
                    content = (message.get("content") or "").strip()
                    if not content:
                        logger.warning("视觉模型返回空 content")
                        return None
                    # 回复前缀过滤：时间戳 [ts]、[私聊 - 名字]、[群聊 - 名字]
                    # （模型偶发把注入的历史/装饰前缀复读进回复时去除）
                    content = strip_reply_prefix(content)
                    return {"kind": "text", "content": content}
                # 任务/提醒工具调用优先原样返回（既有单次语义不变，不让查询
                # 工具往返延迟任务投递）；查询型工具（搜索/抓取/记忆检索）
                # 回填结果继续循环
                other = [tc for tc in tool_calls
                         if _tool_call_name(tc) not in
                         ("web_search", "web_fetch", "recall_memory")]
                if other:
                    fn = other[0].get("function", {}) or {}
                    return {"kind": "tool_call",
                            "name": fn.get("name", ""),
                            "arguments": fn.get("arguments", "")}
                # 回填：assistant(tool_calls) + 每个调用的 tool 结果消息。
                # 就地 append——payload["messages"] 与本地列表同引用，末轮
                # 循环耗尽后直接退出（最后一次要求的工具调用不再执行）
                messages.append({"role": "assistant",
                                 "content": message.get("content") or "",
                                 "tool_calls": tool_calls})
                for tc in tool_calls:
                    name = _tool_call_name(tc)
                    if name == "web_search":
                        query = _tool_arg(tc, "query")
                        if not query:
                            tool_text = "web_search 参数缺少 query"
                        else:
                            try:
                                tool_text = format_search_results(
                                    query, web_search(query))
                            except WebSearchError as e:
                                tool_text = f"搜索暂时不可用：{e}"
                        logger.info(f"[搜索] {query!r} -> {tool_text[:80]}")
                    elif name == "recall_memory":
                        query = _tool_arg(tc, "query")
                        if not chat_id:
                            tool_text = "没有可检索的记忆上下文"
                        elif not query:
                            tool_text = "recall_memory 参数缺少 query"
                        else:
                            try:
                                tool_text = self._recall_memory(chat_id, query)
                            except Exception as e:
                                tool_text = f"记忆检索失败：{e}"
                        logger.info(f"[记忆检索] {query!r} -> {tool_text[:60]}")
                    else:  # web_fetch（任务/提醒已在 other 分支返回）
                        url = _tool_arg(tc, "url")
                        if not url:
                            tool_text = "web_fetch 参数缺少 url"
                        else:
                            try:
                                tool_text = web_fetch(url)
                            except Exception as e:
                                tool_text = f"网页抓取失败：{e}"
                        logger.info(f"[抓取] {url[:60]} -> {tool_text[:60]}")
                    messages.append({"role": "tool",
                                     "tool_call_id": tc.get("id") or "",
                                     "content": tool_text})
            logger.warning("[vision] web_search 循环达补全次数上限，放弃")
            return None
        except Exception as e:
            logger.error(f"视觉模型调用失败: {e}")
            return None

    def _route_vision_result(self, chat_name, sender, result, img_paths=None):
        """vision 单调用结果分流 hook（可覆写）。

        result 为 call_vision_api 的 dict 返回（{'kind':'tool_call',
        'name','arguments'} | {'kind':'text','content'}）或 None；基类默认
        返回 None（未处理），由 AgentBot（xiaoli_bot）覆写：tool_call →
        任务投递，text → 直接回复。img_paths 为截图临时文件路径列表（调用
        方 finally 负责清理，覆写方需在返回前同步消费）。"""
        return None

    # 视觉模型输入最长边上限：屏幕截图（可能 4K 全窗口）base64 直发体积过大
    MAX_IMAGE_EDGE = 2048

    def _save_screenshot_compressed(self, image, path):
        """缩放 + JPEG 压缩保存截图（pillow 已在依赖）。返回文件字节数。
        最长边超 MAX_IMAGE_EDGE 时等比缩放到上限内；PIL 不可用/失败时
        退回原样 PNG 保存（保证图片处理链路不因压缩失败中断）。"""
        try:
            from PIL import Image
            img = image.convert("RGB")
            w, h = img.size
            longest = max(w, h)
            if longest > self.MAX_IMAGE_EDGE:
                scale = self.MAX_IMAGE_EDGE / longest
                img = img.resize((max(1, int(w * scale)), max(1, int(h * scale))),
                                 Image.LANCZOS)
            img.save(path, format="JPEG", quality=88, optimize=True)
        except Exception as e:
            logger.error(f"[处理] 图片压缩失败，退回原样保存: {e}")
            try:
                image.save(path, format="PNG")
            except Exception as e2:
                logger.error(f"[处理] 退回保存也失败: {e2}")
        return os.path.getsize(path) if os.path.isfile(path) else 0

    def _capture_media_images(self, chat_name, min_top=None, exclude_rows=None):
        """捕获消息区对方本轮全部新媒体，返回临时文件路径列表（时间正序）。

        min_top：消息区 1x 下沿阈值（上层传 analyze_window 的 bot_bottom）
        ——只捕获 top ≥ 阈值的媒体框，排除 bot 自己的文件卡片与历史媒体；
        None = 不过滤（全量）。
        exclude_rows：消息区 1x 坐标 y 区间 [(top, bottom), ...]（上层传文件
        名 OCR 行区间）——与文件行垂直相交的媒体框剔除（文件卡片的类型
        图标会从面板跳出成独立小媒体框，面板本身判气泡）。

        每张独立走「点击 → 查看器判定分支」——微信 PC 每条消息各带一个
        头像，多张图片是多个独立媒体框（连通域按背景缝隙切分，不会被粘连）：
        - 真图片：点击打开「图片和视频」查看器 → Ctrl+C 复制原图进剪贴板
          （CF_HDROP 原始分辨率，替代旧预览窗截屏——截屏受窗口尺寸/DPI
          限制且被遮挡会截到遮挡物）→ ESC 关查看器。ESC 只在确认查看器
          存在时才按（表情路径按 ESC 会关掉微信主窗口，真机事故）。剪贴板
          为空（复制失败）→ 落表情路线兜底（先关查看器，避免遮挡裁剪区）
        - 没开（表情包）：全程不碰 ESC/Ctrl+C，截微信主窗口按媒体矩形
          裁剪表情本体送视觉模型（动图取当前帧）
        单张失败（复制为空/裁剪越界/异常）跳过该张，不拖垮整批。

        微信 RWTemp 临时文件生命周期不受控，复制一份到自己的临时文件再返回。
        """
        boxes_fn = getattr(self.wx, "media_screen_boxes", None)
        rects = (boxes_fn(min_top=min_top, exclude_rows=exclude_rows)
                 if boxes_fn is not None else [])
        paths = []
        for (ml, mt, mr, mb) in rects:
            try:
                pyautogui.click((ml + mr) // 2, (mt + mb) // 2)
                time.sleep(1.0)  # 等预览窗打开
                viewer_open = find_window_by_title("图片和视频") is not None
                copied = self._copy_image_from_viewer() if viewer_open else None
                if viewer_open:
                    pyautogui.press('esc')  # 关查看器（确认存在才按，全库唯一 ESC 点）
                    time.sleep(0.3)
                if copied:
                    paths.append(copied)
                    continue
                if viewer_open:
                    logger.warning("[图片复制] 剪贴板复制失败，落表情路线兜底")
                p = self._crop_media_region(ml, mt, mr, mb)
                if p:
                    paths.append(p)
            except Exception as e:
                logger.error(f"[图片捕获] 单张失败，跳过: {e}")
        return paths

    def _copy_image_from_viewer(self):
        """查看器已打开时：Ctrl+C 复制原图 → 读剪贴板 → 返回临时文件路径。

        不负责开关查看器（调用方确认存在并统一 ESC 关闭）；返回 None =
        复制失败，调用方落表情路线。"""
        tmp_path = None
        try:
            pyautogui.hotkey('ctrl', 'c')
            time.sleep(0.5)
            from PIL import ImageGrab
            grabbed = ImageGrab.grabclipboard()
            if isinstance(grabbed, (list, tuple)):
                files = [x for x in grabbed
                         if isinstance(x, str) and os.path.isfile(x)]
                if not files:
                    logger.warning("[图片复制] 剪贴板无有效文件路径")
                    return None
                import shutil
                with tempfile.NamedTemporaryFile(suffix='.jpg', delete=False) as tmpfile:
                    tmp_path = tmpfile.name
                shutil.copyfile(files[0], tmp_path)
                return tmp_path
            if grabbed is not None:
                # 兜底：剪贴板直接是位图（非文件路径），压缩保存
                with tempfile.NamedTemporaryFile(suffix='.jpg', delete=False) as tmpfile:
                    tmp_path = tmpfile.name
                self._save_screenshot_compressed(grabbed, tmp_path)
                return tmp_path
            logger.warning("[图片复制] 剪贴板为空（可能未复制成功）")
            return None
        except Exception as e:
            logger.error(f"[图片复制] 异常: {e}")
            if tmp_path and os.path.exists(tmp_path):
                try:
                    os.unlink(tmp_path)
                except Exception:
                    pass
            return None

    def _crop_media_region(self, ml, mt, mr, mb):
        """表情路线：截微信主窗口并按媒体矩形裁剪表情本体，返回临时文件路径。

        点击无查看器的媒体 = 表情包——截主窗口当前画面裁出该块（动图取
        当前帧）。矩形是点击前测得的屏幕坐标，表情点击不改变布局，仍有效。"""
        try:
            hwnd = find_window_by_title("微信")
            if not hwnd:
                logger.warning("[表情] 未找到微信主窗口")
                return None
            ensure_window_visible(hwnd)
            wr = window_rect(hwnd)
            if not wr:
                return None
            shot = pyautogui.screenshot(region=wr)
            # 屏幕坐标平移到主窗口图内坐标并夹紧边界（表情不可能越界，
            # 夹紧只是防测量瞬间窗口移动的脏数据）
            l = max(0, ml - wr[0])
            t = max(0, mt - wr[1])
            r = min(shot.width, mr - wr[0])
            b = min(shot.height, mb - wr[1])
            if r - l < 10 or b - t < 10:
                logger.warning("[表情] 媒体矩形越界，放弃裁剪")
                return None
            crop = shot.crop((l, t, r, b))
            with tempfile.NamedTemporaryFile(suffix='.jpg', delete=False) as tmpfile:
                tmp_path = tmpfile.name
            self._save_screenshot_compressed(crop, tmp_path)
            logger.info(f"[表情] 点击无查看器，按表情路线裁剪 ({r - l}x{b - t})")
            return tmp_path
        except Exception as e:
            logger.error(f"[表情] 裁剪异常: {e}")
            return None

    def _process_pure_image(self, chat_name, min_top=None):
        """纯图/表情消息处理：捕获对方本轮全部新媒体 → vision 单调用（全部
        图 + 人设 + 历史，一次看图角色化回复/任务判定）→ dict 原样路由给
        _route_vision_result（不压文本；img_paths 由调用方 finally 清理，
        覆写方需在返回前同步消费）。sender 无独立来源，以 chat_name 兜底。
        与图+文路径（_vision_route）同一链路语义：角色化回复而非旧两段式
        的客观图片复述。"""
        tmp_paths = self._capture_media_images(chat_name, min_top=min_top)
        if not tmp_paths:
            return None
        content = [{"type": "text", "text": VISION_IMAGE_PROMPT}]
        try:
            for p in tmp_paths:
                with open(p, 'rb') as f:
                    img_bytes = f.read()
                content.append({
                    "type": "image_url",
                    "image_url": {
                        "url": f"data:image/png;base64,{base64.b64encode(img_bytes).decode()}"}},
                )
            result = self.call_vision_api(content, chat_id=chat_name)
            if result:
                return self._route_vision_result(chat_name, chat_name, result,
                                                 img_paths=tmp_paths)
            return None
        except Exception as e:
            logger.error(f"[图片处理] 异常: {e}")
            return None
        finally:
            for p in tmp_paths:
                try:
                    os.unlink(p)
                except Exception:
                    pass

    def _extract_file_text(self, filepath):
        """从文件中提取文本内容，支持纯文本、docx/doc（Office COM）、xlsx/xls。
        PDF 暂不支持（历史上此处声明过但从未实现，PDF 会走文本读取失败返回 None）。"""
        filename = os.path.basename(filepath)
        ext = os.path.splitext(filepath)[1].lower()

        # 纯文本文件
        text_extensions = {
            '.txt', '.py', '.java', '.js', '.ts', '.html', '.css', '.json',
            '.xml', '.yaml', '.yml', '.md', '.csv', '.log', '.ini', '.cfg',
            '.sh', '.bat', '.c', '.cpp', '.h', '.hpp', '.rs', '.go', '.rb',
            '.php', '.sql', '.r', '.m', '.swift', '.kt', '.scala', '.lua',
            '.toml', '.tex', '.svg', '.pl', '.ps1', '.conf', '.properties',
        }
        if ext in text_extensions:
            for enc in ('utf-8', 'gbk', 'gb2312', 'latin-1'):
                try:
                    with open(filepath, 'r', encoding=enc) as f:
                        return f.read()
                except (UnicodeDecodeError, UnicodeError):
                    continue
            logger.warning(f"[文件] 无法以任何编码读取: {filename}")
            return None

        # .docx
        if ext == '.docx':
            try:
                import docx
                doc = docx.Document(filepath)
                text = '\n'.join([para.text for para in doc.paragraphs])
                return text if text.strip() else None
            except ImportError:
                logger.warning("[文件] 未安装 python-docx 库")
                return None
            except Exception as e:
                logger.warning(f"[文件] 读取 docx 失败: {e}")
                return None

        # .doc（旧版 Word）
        if ext == '.doc':
            text = self._extract_office_com_text(filepath, 'Word.Application')
            if text:
                return text
            return None

        # .pptx
        if ext == '.pptx':
            try:
                import pptx
                prs = pptx.Presentation(filepath)
                slides_text = []
                for i, slide in enumerate(prs.slides, 1):
                    slide_lines = [f"--- 幻灯片 {i} ---"]
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            for para in shape.text_frame.paragraphs:
                                if para.text.strip():
                                    slide_lines.append(para.text)
                    if len(slide_lines) > 1:
                        slides_text.append('\n'.join(slide_lines))
                result = '\n\n'.join(slides_text)
                return result if result.strip() else None
            except ImportError:
                logger.warning("[文件] 未安装 python-pptx 库")
                return None
            except Exception as e:
                logger.warning(f"[文件] 读取 pptx 失败: {e}")
                return None

        # .ppt（旧版 PowerPoint）
        if ext == '.ppt':
            text = self._extract_office_com_text(filepath, 'PowerPoint.Application')
            if text:
                return text
            return None

        # .xlsx
        if ext == '.xlsx':
            try:
                import openpyxl
                wb = openpyxl.load_workbook(filepath, read_only=True, data_only=True)
                all_text = []
                for sheet_name in wb.sheetnames:
                    ws = wb[sheet_name]
                    sheet_lines = [f"--- 工作表: {sheet_name} ---"]
                    for row in ws.iter_rows(values_only=True):
                        row_text = '\t'.join([
                            str(cell) if cell is not None else '' for cell in row
                        ])
                        if row_text.strip():
                            sheet_lines.append(row_text)
                    all_text.append('\n'.join(sheet_lines))
                wb.close()
                result = '\n\n'.join(all_text)
                return result if result.strip() else None
            except ImportError:
                logger.warning("[文件] 未安装 openpyxl 库")
                return None
            except Exception as e:
                logger.warning(f"[文件] 读取 xlsx 失败: {e}")
                return None

        # .xls（旧版 Excel）
        if ext == '.xls':
            # 优先用 xlrd，失败了用 Excel COM
            try:
                import xlrd
                wb = xlrd.open_workbook(filepath)
                all_text = []
                for sheet in wb.sheets():
                    sheet_lines = [f"--- 工作表: {sheet.name} ---"]
                    for row_idx in range(sheet.nrows):
                        row_values = sheet.row_values(row_idx)
                        row_text = '\t'.join([
                            str(cell) if cell != '' else '' for cell in row_values
                        ])
                        if row_text.strip():
                            sheet_lines.append(row_text)
                    all_text.append('\n'.join(sheet_lines))
                result = '\n\n'.join(all_text)
                return result if result.strip() else None
            except ImportError:
                pass
            except Exception as e:
                logger.debug(f"[文件] xlrd 读取失败: {e}")
            # 回退到 Excel COM
            text = self._extract_office_com_text(filepath, 'Excel.Application')
            if text:
                return text
            return None

        # 未知扩展名（含 PDF，暂不支持解析）尝试按文本读取；二进制读取失败返回 None
        try:
            with open(filepath, 'r', encoding='utf-8') as f:
                text = f.read()
            if text.strip():
                logger.debug(f"[文件] 未知扩展名 .{ext}，按文本读取成功")
                return text
        except Exception:
            pass

        return None

    def _reply_with_file(self, chat_name, sender, file_text, filename):
        """根据文件内容生成回复并发送"""
        logger.debug(f"[文件处理] 内容前100字符: {file_text[:100]}...")
        self._add_history(chat_name, "assistant",
                          f"[文件内容: {filename}] {file_text}")
        refine_prompt = (
            f"用户发来一个文件（{filename}），内容如下：\n\n"
            f"{file_text}\n\n"
            f"请根据这个文件内容，以{self.nickname}的身份回复用户。"
        )
        is_group = is_group_chat(chat_name)
        final_reply = self.call_chat_ai(chat_name, refine_prompt,
                                        sender_name=sender, is_group=is_group)
        self._send_text(final_reply, chat_name)
        return True

    def _extract_file_display_name(self, msg):
        """从 FileMessage 提取显示文件名。
        wxauto4 的 content 格式：'文件\\n<文件名>\\n[<大小>\\n]微信电脑版'
        （实测：'文件\\n养生规划表.html\\n微信电脑版'）
        返回文件名或 None"""
        try:
            content = getattr(msg, "content", "") or ""
            m = re.search(r"^文件\n([^\n]+)", content, flags=re.M)
            if m:
                return m.group(1).strip()
            # 兜底：按 repattern 解析
            rep = getattr(msg, "repattern", None)
            if rep:
                m2 = re.search(rep, content)
                if m2 and m2.group(1):
                    return m2.group(1).strip()
            # 视觉后端兼容：content 即显示文件名（无 '文件\n' 前缀）
            if content and "\\n" not in content and "\n" not in content:
                # 纯文件名（不含换行/前缀）——视觉后端 file 消息格式。
                # 但消息区 OCR 可能把多条文件消息合并成一个文本块
                # （实测 '新宣传.docx 部门简介+纳新宣传.docx W'，文件图标被
                # OCR 成尾部杂字符）——整串当文件名必然匹配失败，需先拆出
                # 真实文件名（含常见文档扩展名的 token，取第一个）。
                name = content.strip()
                if name:
                    return _extract_file_name_token(name) or name
        except Exception as e:
            logger.error(f"[文件] 提取文件名失败: {e}")
        return None

    def _find_file_by_display_name(self, display_name):
        """按消息中的显示文件名在接收目录定位对方发来的文件。

        实现在模块级 _find_file_by_display_name_impl（纯函数，单测直测）：
        分隔符归一化包含匹配 + (N) 重名编号最大优先 + ctime 平局。"""
        return _find_file_by_display_name_impl(self.file_storage_path,
                                               display_name)

    def _process_file(self, chat_name, sender, msg_obj):
        """处理文件消息：按消息显示名在微信接收目录定位（wxauto4 download() 不可用）
        -> 提取文字 -> 发送给 AI"""
        logger.info(f"📁 收到 {sender} 的文件，开始处理...")

        try:
            # 1. 从消息提取显示文件名并定位接收目录中的文件
            display_name = self._extract_file_display_name(msg_obj)
            file_path = self._find_file_by_display_name(display_name) if display_name else None
            if not file_path:
                logger.error(f"[文件] 未在接收目录定位到文件: {display_name!r}")
                self._send_text("文件下载失败，请重试～", chat_name)
                return False

            filename = os.path.basename(file_path)
            file_size = os.path.getsize(file_path)
            logger.info(f"📁 文件已定位: {filename} ({file_size} bytes)")

            # 2. 提取文本
            text_content = self._extract_file_text(file_path)
            if text_content is None:
                logger.warning(f"[文件] 无法提取文本（格式不支持或内容为空）")
                self._send_text(f"收到文件「{filename}」，但这个格式我看不懂呢～", chat_name)
                return False

            logger.info(f"📁 文件「{filename}」提取了 {len(text_content)} 字符，发送给 AI...")
            return self._reply_with_file(chat_name, sender, text_content, filename)

        except Exception as e:
            logger.error(f"[文件] ❌ 异常: {e}", exc_info=True)
            self._send_text("文件处理出错，请重试～", chat_name)
            return False

    def _extract_office_com_text(self, filepath, app_name):
        """通过 Office COM 自动化提取旧格式（.doc/.ppt/.xls）文本，失败则二进制兜底"""
        # 方法1: Office COM 自动化
        try:
            import comtypes.client
            app = comtypes.client.CreateObject(app_name)
            app.Visible = False

            if 'Word' in app_name:
                doc = app.Documents.Open(filepath)
                text = doc.Content.Text
                doc.Close()
            elif 'PowerPoint' in app_name:
                prs = app.Presentations.Open(filepath, WithWindow=False)
                slides = []
                for slide in prs.Slides:
                    for shape in slide.Shapes:
                        if shape.HasTextFrame:
                            slides.append(shape.TextFrame.TextRange.Text)
                text = '\n'.join(slides)
                prs.Close()
            elif 'Excel' in app_name:
                wb = app.Workbooks.Open(filepath)
                sheets = []
                for sheet in wb.Sheets:
                    used = sheet.UsedRange
                    if used:
                        rows = []
                        for row in used.Rows:
                            cells = []
                            for cell in row.Cells:
                                v = cell.Value
                                cells.append(str(v) if v is not None else '')
                            rows.append('\t'.join(cells))
                        sheets.append(
                            f"--- 工作表: {sheet.Name} ---\n" + '\n'.join(rows)
                        )
                text = '\n\n'.join(sheets)
                wb.Close()
            else:
                app.Quit()
                return None

            app.Quit()
            if text and text.strip():
                logger.debug(f"[文件] {app_name} COM 提取成功 ({len(text)} 字符)")
                return text.strip()
        except Exception as e:
            logger.debug(f"[文件] {app_name} COM 失败: {e}")

        # 方法2: 从二进制中提取可读文本（兜底方案）
        try:
            with open(filepath, 'rb') as f:
                data = f.read()
            text_parts = []
            buf = []
            for byte in data:
                if 32 <= byte < 127 or byte in (9, 10, 13):
                    buf.append(chr(byte))
                else:
                    if len(buf) >= 4:
                        text_parts.append(''.join(buf))
                    buf = []
            if len(buf) >= 4:
                text_parts.append(''.join(buf))
            result = '\n'.join(text_parts)
            if len(result) > 100:
                logger.debug(f"[文件] 二进制提取成功 ({len(result)} 字符)")
                return result
        except Exception as e:
            logger.debug(f"[文件] 二进制提取失败: {e}")

        return None

    def call_chat_ai(self, chat_id, user_msg, sender_name=None, is_group=False, multi_sender=False):
        headers = {
            "Authorization": f"Bearer {self.api_key}",
            "Content-Type": "application/json"
        }
        if is_group:
            # 群聊格式 = 群聊名 + 发送者名 + 内容（用户原话：群聊：XXX XXX：消息内容）。
            # 显式三分支：
            #   1. multi_sender=True：user_msg 已由 _handle_unread_session 合并成
            #      「发送者A：内容A\n发送者B：内容B」（每条自带发送者名），
            #      只包『群聊：群名』前缀，不再重包 sender（否则双层嵌套）。
            #   2. sender_name 存在：群聊名 + 发送者名 + 内容（单条，现状）。
            #   3. sender_name 缺失（极端 OCR 失败）：群聊名兜底并打日志，
            #      绝不落入「群聊：{user_msg}」无名字退化分支。
            # multi_sender 是唯一显式信号——禁止用 user_msg 含换行等隐式判断
            # （单条多行文本消息会误判）。
            if multi_sender:
                decorated = f"群聊：{chat_id} {user_msg}"
            elif sender_name:
                decorated = f"群聊：{chat_id} {sender_name}：{user_msg}"
            else:
                logger.warning(f"[群聊] {chat_id} 视觉层未读到发送者名，用群聊名兜底")
                decorated = f"群聊：{chat_id}：{user_msg}"
        else:
            decorated = f"私聊 - {sender_name}：{user_msg}" if sender_name else f"私聊：{user_msg}"
        current_time = time.strftime("%Y-%m-%d %H:%M:%S")
        # 消息布局（缓存友好，与 call_vision_api 对齐）：稳定前缀在前
        # （人设 → 重要记忆 → 历史），每轮变化区在尾（相关记忆 → 当前时间
        # → 当前消息）。「当前时间」每秒变化，绝不能插在历史之前打断前缀。
        messages = [
            {"role": "system", "content": self.system_prompt},
        ]
        important = self._important_block(chat_id)
        if important:
            messages.append({"role": "system", "content": important})
        for h in self._get_history(chat_id):
            if "time" in h:
                ts = h["time"]
                msg_content = f"[{ts}] {h['content']}"
            else:
                msg_content = h['content']
            messages.append({"role": h["role"], "content": msg_content})
        related = self._match_related_memory(chat_id, user_msg)
        if related:
            messages.append({"role": "system", "content": related})
        messages.append({"role": "system", "content": f"当前时间：{current_time}"})
        messages.append({"role": "user", "content": decorated})

        with self._model_lock:
            # 空模型兜底（活跃卡缺失被模板补建后 chat_model 为空）：空串原样
            # 发出必 400，与 vision 链路共用同一兜底纯名
            model = self.chat_model or VISION_MODEL_DEFAULT
            temp = self.chat_temperature
            top_p = self.chat_top_p

        # 上下文预算裁剪：文件全文/超长历史会撑爆模型上下文上限
        # （实测请求 272 万 token → API 400 "maximum context length"）。
        # 从最旧历史开始丢弃，保证单次请求不超模型上下文。
        messages = fit_messages_in_budget(
            messages, budget=getattr(self, "max_context_tokens", 100000))

        payload = {
            "model": model,
            "messages": messages,
            "temperature": temp,
            "top_p": top_p
        }
        try:
            data = self._post_chat_completions(
                self.api_url, headers, payload, self.api_timeout, label="chat",
                meta={"kind": "chat", "model": model, "messages": messages})
        except ApiCallError as e:
            logger.error(f"聊天 API 调用失败: {e}")
            return random.choice(FRIENDLY_API_ERROR_REPLIES)
        choices = data.get("choices", [])
        if choices and "message" in choices[0]:
            reply = choices[0]["message"]["content"]
        else:
            logger.error(f"API返回异常 choices 为空: {data}")
            return random.choice(FRIENDLY_API_ERROR_REPLIES)
        reply = strip_reply_prefix(reply)
        self._add_history(chat_id, "user", decorated)
        self._add_history(chat_id, "assistant", reply)
        return reply

    def _send_text(self, text, chat, placeholder=False):
        """发送文本到指定聊天。placeholder=True 表示这是"处理中"占位消息
        （永不拆分、计数 +1）；默认 False（普通回复）按"至少一个空行"拆分
        成多条发送，并把该聊天的占位计数归零。
        计数按 chat_name 隔离（_pending_placeholders），非全局。"""
        try:
            # 占位消息永不拆分，单条发送，计数语义不变
            if placeholder:
                cleaned_text = text.strip()
                self.wx.send_text(chat, cleaned_text)
                preview = cleaned_text[:50].replace('\n', ' ')
                logger.info(f"🤖 → [{chat}]: {preview}")
                self._pending_placeholders[chat] = self._pending_placeholders.get(chat, 0) + 1
                return

            # 按换行拆分（用户定案：单换行也分次发送——人设要求碎句多段，
            # 模型用单 \n 分句时旧逻辑会整段一起发）。任意连续换行都拆，
            # 空段丢弃；拆不出多段时走原单条发送
            parts = [p.strip() for p in re.split(r'\n+', text) if p.strip()]
            if len(parts) <= 1:
                cleaned_text = text.strip()
                self.wx.send_text(chat, cleaned_text)
                preview = cleaned_text[:50].replace('\n', ' ')
                logger.info(f"🤖 → [{chat}]: {preview}")
            else:
                for part in parts:
                    self.wx.send_text(chat, part)
                    preview = part[:50].replace('\n', ' ')
                    logger.info(f"🤖 → [{chat}]: {preview}")

            # 实质回复归零占位，只执行一次（pop 对未占位过的 chat 安全）
            self._pending_placeholders.pop(chat, None)
        except Exception as e:
            logger.error(f"发送失败: {e}")

    def fetch_models(self):
        url = models_endpoint(self.api_url)
        headers = {"Authorization": f"Bearer {self.api_key}"}
        try:
            resp = requests.get(url, headers=headers, timeout=10)
            if resp.status_code == 200:
                return [m["id"] for m in resp.json().get("data", [])]
            else:
                logger.warning(f"获取模型列表失败: {resp.status_code}")
                return None
        except Exception as e:
            logger.error(f"请求模型列表异常: {e}")
            return None

    def process_new_messages(self):
        self._flush_memory_if_due()
        if self.paused:
            return
        if time.time() - self.last_reply_time < self.cooldown:
            return
        try:
            # 红圈驱动（visual 后端）；降级（wxauto 等）走全量会话遍历
            if hasattr(self.wx, "iter_unread_sessions"):
                sessions = list(self.wx.iter_unread_sessions())
            else:
                sessions = list(self.wx.iter_sessions())
            if not sessions:
                return
            file_re = re.compile(
                r"\.(?:docx?|xlsx?|pptx?|pdf|txt|md|html?|json|csv|zip|rar|7z|png|jpe?g|gif|mp4|mp3)\b",
                re.I)
            for chat_name in sessions:
                if not chat_name:
                    continue
                # 窗口边界：定位 bot 最后回复之后的对方消息（visual 后端）
                analyze = getattr(self.wx, "analyze_window", None)
                win = analyze(chat_name) if analyze else None
                bot_bottom = win.get("bot_bottom") if win else None
                if win is not None:
                    if not (win.get("has_other") or win.get("has_text") or win.get("has_media")):
                        continue
                    if win.get("has_media"):
                        time.sleep(10)  # 防对方话没说完
                        win = analyze(chat_name)
                        if not (win.get("has_other") or win.get("has_text") or win.get("has_media")):
                            continue
                        bot_bottom = win.get("bot_bottom")
                msgs = self.wx.get_messages(chat_name, assume_switched=True)
                window_msgs = [
                    m for m in msgs
                    if m.sender not in (None, "self", self.nickname)
                    and (bot_bottom is None or (m.y is not None and m.y >= bot_bottom))
                ]
                file_text = next(
                    (m.content.strip() for m in window_msgs if file_re.search(m.content or "")),
                    None)
                text_candidates = [
                    m for m in window_msgs
                    if m.content.strip() and not file_re.search(m.content or "")
                ]
                if len(text_candidates) > 1:
                    # 多发送者合并：每条带各自发送者名（不整批只带最后一条）
                    multi_sender = True
                    text_parts = [
                        f"{m.sender or chat_name}：{m.content.strip()}"
                        for m in text_candidates
                    ]
                else:
                    multi_sender = False
                    text_parts = [m.content.strip() for m in text_candidates]
                text_content = "\n".join(text_parts)
                has_media = bool(win.get("has_media")) if win else False
                sender = window_msgs[-1].sender if window_msgs else chat_name
                # 分发
                if file_text:
                    logger.info(f"📁 判断为文件消息：{chat_name}")
                    self._process_file(chat_name, sender, window_msgs[-1])
                    self.last_reply_time = time.time()
                    return
                if has_media and not text_content:
                    logger.info(f"🖼 判断为图片消息：{chat_name}")
                    if not self._process_pure_image(chat_name):
                        # 失败必须回一句：发送点输入框顺带清红圈，防滞留循环
                        self._send_text("图片识别失败了，可能是什么地方出了问题呀～",
                                        chat_name)
                    self.last_reply_time = time.time()
                    return
                if text_content:
                    is_group = is_group_chat(chat_name)
                    question = text_content
                    if is_group:
                        at_tag = f"@{self.nickname}"
                        if at_tag not in question:
                            continue
                        question = question.replace(at_tag, "").strip()
                        if not question:
                            question = "你好呀～"
                    logger.info(f"💬 [{chat_name}] {sender}: {question[:80]}")
                    reply = self.call_chat_ai(chat_name, question, sender_name=sender, is_group=is_group, multi_sender=multi_sender)
                    self._send_text(reply, chat_name)
                    self.last_reply_time = time.time()
                    return
        except Exception as e:
            logger.error(f"处理消息异常: {e}\n{traceback.format_exc()}")

    def run(self, stop_event=None, poll_interval=0.5):
        state = "暂停中，输入 resume 开始回复" if self.paused else "运行中"
        logger.info(f"✅ 小漓已启动（{state}）")
        while True:
            if stop_event is not None and stop_event.is_set():
                logger.info("🛑 引擎已停止")
                self._flush_memory()  # 退出前落盘节流窗口内的记忆
                return
            try:
                self.process_new_messages()
            except Exception as e:
                logger.error(f"主循环异常: {e}\n{traceback.format_exc()}")
            # 可中断睡眠：stop 响应延迟 ≤100ms（控制台模式不传 stop_event，行为不变）
            # poll_interval 0.5s 恒定快档（与 GUI 引擎一致，不间断监听）
            remain = poll_interval
            while remain > 1e-9:
                if stop_event is not None and stop_event.is_set():
                    logger.info("🛑 引擎已停止")
                    self._flush_memory()
                    return
                time.sleep(min(0.1, remain))
                remain -= 0.1


class Controller:
    """命令行控制器：命令注册表分发（子类可扩展命令）。

    设计：命令首词 → handler（接收完整命令串含参数）。新增命令 = 注册
    dict 加一项 + 实现 handler，不再复制整段 if/elif 链（历史：Tianshu
    Controller 复制 _listen 150 行只为了加两个命令）。
    quit 语义：置 stop_event + flush 记忆 → bot.run(stop_event) 退出 →
    主线程正常收尾（不再 os._exit 硬退，避免丢节流窗口内的记忆）。"""

    # 命令名 → 一行说明（help 命令展示；子类合并扩展）
    HELP = {
        "pause": "暂停自动回复",
        "resume": "恢复自动回复",
        "model": "切换聊天模型（model 交互选择 / model <名称> 直接切换并持久化）",
        "chat-temp": "设置聊天模型温度 (0~2)",
        "chat-top-p": "设置聊天模型 top_p (0~1)",
        "clear": "清空对话历史（全部 / clear <聊天ID>）",
        "del": "删除聊天中的消息：del <聊天ID> <序号1> [序号2] [3-5]",
        "memory": "查看聊天历史消息（带序号）",
        "status": "查看当前状态（模型信息、温度、top_p）",
        "quit": "退出程序",
        "help": "显示本帮助",
    }

    def __init__(self, bot):
        self.bot = bot
        self.stop_event = threading.Event()  # quit 置位 → bot.run 退出（优雅收尾）
        self._commands = {}
        self._register_commands()

    def _register_commands(self):
        """注册命令处理器（key = 命令首词，handler 接收完整命令串）。"""
        self._commands.update({
            "help": self._cmd_help,
            "pause": self._cmd_pause,
            "resume": self._cmd_resume,
            "model": self._cmd_model,
            "chat-temp": self._cmd_chat_temp,
            "chat-top-p": self._cmd_chat_top_p,
            "clear": self._cmd_clear,
            "del": self._cmd_del,
            "memory": self._cmd_memory,
            "status": self._cmd_status,
            "quit": self._cmd_quit,
        })

    def start(self):
        threading.Thread(target=self._listen, daemon=True).start()

    def _listen(self):
        while True:
            try:
                cmd = input().strip().lower()
                if not cmd:
                    continue
                name = cmd.split(" ", 1)[0]
                handler = self._commands.get(name)
                if handler is None:
                    logger.info(f"❌ 未知命令: {cmd}")
                    continue
                handler(cmd)
            except (EOFError, KeyboardInterrupt):
                break

    # ---------- 命令处理器 ----------

    def _cmd_help(self, cmd):
        lines = ["可用命令："]
        for name, handler in self._commands.items():
            lines.append(f"  {name:<12} - {self.HELP.get(name, '')}")
        print("\n".join(lines))

    def _cmd_pause(self, cmd):
        self.bot.paused = True
        logger.info("⏸️  已暂停自动回复")

    def _cmd_resume(self, cmd):
        self.bot.paused = False
        logger.info("▶️  已恢复自动回复")

    def _cmd_model(self, cmd):
        rest = cmd[len("model"):].strip()
        if not rest:
            self._select_model("chat")
            return
        with self.bot._model_lock:
            self.bot.chat_model = rest
        logger.info(f"🔄 聊天模型已切换为：{rest}")
        self._persist_model_setting("chat_model", rest)

    def _persist_model_setting(self, key, value):
        """CLI 切模型持久化：写回 config.json + 活跃角色卡（GUI 模式下
        投影字段以卡为准，写卡才能跨重启生效）。失败仅告警不中断。"""
        try:
            with open("config.json", "r", encoding="utf-8") as f:
                cfg = json.load(f)
            cfg[key] = value
            with open("config.json", "w", encoding="utf-8") as f:
                json.dump(cfg, f, ensure_ascii=False, indent=4)
            try:
                from xiaoli_app import card_store
                cid = cfg.get("active_card_id")
                if cid:
                    cards_dir = os.path.join(
                        os.path.dirname(os.path.abspath("config.json")), "cards")
                    card = card_store.get_card(cards_dir, cid)
                    if card is not None:
                        card[key] = value
                        card_store.save_card(cards_dir, card)
            except Exception:
                pass
        except Exception as e:
            logger.warning(f"保存配置失败: {e}")

    def _cmd_chat_temp(self, cmd):
        self._set_numeric(cmd, "chat-temp", 0, 2, "温度", self.bot.set_chat_temperature)

    def _cmd_chat_top_p(self, cmd):
        self._set_numeric(cmd, "chat-top-p", 0, 1, "top_p", self.bot.set_chat_top_p)

    def _set_numeric(self, cmd, name, lo, hi, label, setter):
        rest = cmd[len(name):].strip()
        if not rest:
            logger.warning(f"用法: {name} <值>（{lo}~{hi}）")
            return
        try:
            val = float(rest)
            if lo <= val <= hi:
                setter(val)
            else:
                logger.warning(f"{label}值应在 {lo}~{hi} 之间")
        except ValueError:
            logger.warning("请输入有效的数字")

    def _cmd_clear(self, cmd):
        rest = cmd[len("clear"):].strip()
        if rest:
            self.bot.clear_history(rest)
        else:
            self.bot.clear_history()

    def _cmd_del(self, cmd):
        parts = cmd.split()
        if len(parts) < 3:
            logger.info("用法: del <聊天ID> <序号1> [序号2] ... 或 del <聊天ID> <起始序号-结束序号>")
            return
        chat_id = parts[1]
        indices = []
        for part in parts[2:]:
            if '-' in part:
                try:
                    start, end = map(int, part.split('-'))
                    if start > end:
                        start, end = end, start
                    indices.extend(range(start, end + 1))
                except ValueError:
                    logger.warning(f"无效的范围格式: {part}")
            else:
                try:
                    indices.append(int(part))
                except ValueError:
                    logger.warning(f"无效的序号: {part}")
        if not indices:
            return
        print(f"即将从聊天 '{chat_id}' 中删除序号: {sorted(set(indices))}")
        confirm = input("确认删除？(y/n): ").strip().lower()
        if confirm == 'y':
            self.bot.delete_messages(chat_id, indices)
        else:
            logger.info("已取消删除")

    def _cmd_memory(self, cmd):
        target = cmd[len("memory"):].strip()
        self._show_memory(target)

    def _cmd_status(self, cmd):
        self._show_status()

    def _cmd_quit(self, cmd):
        logger.info("👋 程序退出")
        self.stop_event.set()
        self.bot._flush_memory()  # 节流窗口内的记忆落盘（run 退出路径也会 flush，双保险）

    def _show_status(self):
        total_msgs = sum(
            len(v.get("recent", [])) if isinstance(v, dict) else len(v)
            for v in self.bot.memory_db.values())
        deep_msgs = sum(self.bot._deep_count.values())
        chat_count = len(self.bot.memory_db)
        print(f"当前聊天模型：{self.bot.chat_model}")
        print(f"  聊天温度: {self.bot.chat_temperature}")
        print(f"  聊天 top_p: {self.bot.chat_top_p}")
        print(f"视觉模型：随聊天模型（{self.bot.chat_model}，单模型化）")
        print(f"对话记忆：共 {chat_count} 个聊天，近期 {total_msgs} 条"
              f"（深层存档 {deep_msgs} 条）")
        if chat_count > 0:
            print("具体聊天对象：")
            for chat in self.bot.memory_db.keys():
                print(f"  - {chat}")
        print(f"当前状态：{'运行中' if not self.bot.paused else '已暂停'}")

    def _show_memory(self, target):
        matches = [chat for chat in self.bot.memory_db if target in chat]
        if not matches:
            logger.info(f"❌ 未找到包含 '{target}' 的聊天记录")
            return
        chat_key = matches[0]
        st = self.bot.memory_db[chat_key]
        recent = st.get("recent", []) if isinstance(st, dict) else st
        print(f"📂 聊天对象: {chat_key}，近期 {len(recent)} 条消息：")
        for i, msg in enumerate(recent, 1):
            role = "用户" if msg["role"] == "user" else "小漓"
            ts = msg.get("time", "未知时间")
            content = msg["content"]
            print(f"  {i}. [{ts}] {role}: {content}")
        if isinstance(st, dict):
            imp = [x.get("content") for x in st.get("important", [])
                   if isinstance(x, dict) and x.get("content")]
            if imp:
                print("📌 重要记忆：")
                for c in imp:
                    print(f"  - {c}")
            if st.get("index"):
                print(f"🗂 关键词索引：{len(st['index'])} 条"
                      "（命中时自动注入，memory <聊天ID> 不逐条展示）")

    def _select_model(self, model_type):
        was_paused = self.bot.paused
        self.bot.paused = True
        logger.info(f"⏸️  已暂停，正在获取模型列表...")
        models = self.bot.fetch_models()
        if not models:
            logger.error("无法获取模型列表")
            self.bot.paused = was_paused
            return
        for i, m in enumerate(models):
            print(f"  {i+1}. {m}")
        choice = input("请输入序号或模型名，输入 cancel 取消: ").strip()
        if choice.lower() == "cancel":
            logger.info("已取消")
            self.bot.paused = was_paused
            return
        try:
            idx = int(choice) - 1
            if 0 <= idx < len(models):
                new_model = models[idx]
            else:
                logger.error("序号超出范围")
                self.bot.paused = was_paused
                return
        except ValueError:
            new_model = choice
        with self.bot._model_lock:
            # 单模型化：视觉/分类统一走 chat_model，仅 chat 分支可切换
            self.bot.chat_model = new_model
            logger.info(f"🔄 聊天模型已切换为：{new_model}")
        ask = input("是否清空所有对话历史？(y/n): ").strip().lower()
        if ask == 'y':
            self.bot.clear_history()
        self.bot.paused = was_paused
        if not self.bot.paused:
            logger.info("▶️  已自动恢复回复")


if __name__ == "__main__":
    bot = WeChatBot(load_config())
    controller = Controller(bot)
    controller.start()
    bot.run(stop_event=controller.stop_event)